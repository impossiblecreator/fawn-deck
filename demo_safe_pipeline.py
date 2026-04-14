"""
Demo: Safe PPTX editing pipeline with automatic corruption detection.

Creates a temporary copy of the output deck, then runs 3 sequential edits:
  1. A clean edit (adds a text box) — should PASS validation
  2. A deliberately corrupt edit (zero-width shape) — should FAIL and auto-rollback
  3. Another clean edit after recovery — should PASS, proving rollback was clean

This script does NOT modify any real project files.
"""

import os
import shutil
import tempfile
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

from pptx_safe_ops import safe_edit_slide, validate_deck

WORKSPACE = os.path.dirname(os.path.abspath(__file__))
SOURCE_DECK = os.path.join(WORKSPACE, "output", "deck_final.pptx")


def main():
    # Work on a temporary copy so we don't touch the real deck
    tmp_dir = tempfile.mkdtemp(prefix="safe_pipeline_demo_")
    demo_deck = os.path.join(tmp_dir, "demo_deck.pptx")
    shutil.copy2(SOURCE_DECK, demo_deck)

    print("=" * 60)
    print("  SAFE PIPELINE DEMO — 3 sequential slide edits")
    print("=" * 60)

    # Pre-validate
    print("\n--- Pre-validation ---")
    pre = validate_deck(demo_deck)
    print(f"  {pre.summary()}")

    # ---------------------------------------------------------------
    # Edit 1: Clean edit — add a text box to slide 1
    # ---------------------------------------------------------------
    print("\n--- Edit 1: Add text box to slide 1 (should PASS) ---")

    def edit_add_textbox(slide, prs):
        txBox = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(5), Inches(0.5))
        tf = txBox.text_frame
        tf.text = "Pipeline validation: OK"
        for run in tf.paragraphs[0].runs:
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0x12, 0x3C, 0x33)
            run.font.name = "Haas Grot Text Trial"

    result1 = safe_edit_slide(demo_deck, slide_num=1, edit_fn=edit_add_textbox)
    _print_result("Edit 1", result1)

    # ---------------------------------------------------------------
    # Edit 2: Deliberately corrupt — zero-width text box on slide 3
    # ---------------------------------------------------------------
    print("\n--- Edit 2: Zero-width text box on slide 3 (should FAIL + rollback) ---")

    def edit_zero_width(slide, prs):
        # This creates a text box with zero width — validation will catch it
        txBox = slide.shapes.add_textbox(Inches(2), Inches(2), Emu(0), Inches(1))
        txBox.text_frame.text = "This should never persist"

    result2 = safe_edit_slide(demo_deck, slide_num=3, edit_fn=edit_zero_width)
    _print_result("Edit 2", result2)

    # ---------------------------------------------------------------
    # Edit 3: Clean edit after rollback — add text to slide 5
    # ---------------------------------------------------------------
    print("\n--- Edit 3: Add text to slide 5 after recovery (should PASS) ---")

    def edit_post_recovery(slide, prs):
        txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(6), Inches(0.6))
        tf = txBox.text_frame
        tf.text = "Post-recovery edit: pipeline is healthy"
        for run in tf.paragraphs[0].runs:
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0x56, 0x53, 0x4F)
            run.font.name = "Haas Grot Text Trial"

    result3 = safe_edit_slide(demo_deck, slide_num=5, edit_fn=edit_post_recovery)
    _print_result("Edit 3", result3)

    # Final validation
    print("\n--- Final validation ---")
    final = validate_deck(demo_deck)
    print(f"  {final.summary()}")

    # Summary
    print("\n" + "=" * 60)
    print("  RESULTS SUMMARY")
    print("=" * 60)
    print(f"  Edit 1 (clean):     {'PASS' if result1.success else 'FAIL'}  rolled_back={result1.rolled_back}")
    print(f"  Edit 2 (corrupt):   {'PASS' if result2.success else 'FAIL'}  rolled_back={result2.rolled_back}")
    print(f"  Edit 3 (recovery):  {'PASS' if result3.success else 'FAIL'}  rolled_back={result3.rolled_back}")
    print()

    if result1.success and not result2.success and result2.rolled_back and result3.success:
        print("  All 3 edits behaved correctly.")
        print("  The pipeline caught corruption in edit 2 and auto-recovered.")
    else:
        print("  WARNING: Unexpected results — review output above.")

    # Clean up
    shutil.rmtree(tmp_dir, ignore_errors=True)
    print(f"\n  Temp files cleaned up. No project files were modified.")


def _print_result(label, result):
    status = "PASS" if result.success else "FAIL"
    print(f"  [{status}] {result.message}")
    if result.rolled_back:
        print(f"  -> AUTO-ROLLBACK performed")
    if result.report.errors:
        for e in result.report.errors:
            print(f"     [ERROR] {e}")
    if result.report.warnings:
        for w in result.report.warnings:
            print(f"     [WARN]  {w}")


if __name__ == "__main__":
    main()
