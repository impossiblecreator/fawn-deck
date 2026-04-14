"""Edit slide 4 - The Solution (three features)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Load presentation
prs = Presentation('workers/worker_A.pptx')
slide = prs.slides[3]  # Slide 4 (0-indexed)

# Brand colors
SOFT_IVORY = RGBColor(246, 241, 233)
PHTHALO_GREEN = RGBColor(18, 60, 51)
WARM_BEIGE = RGBColor(231, 216, 199)
MUSHROOM_TAUPE = RGBColor(184, 169, 153)
BODY_GRAY = RGBColor(86, 83, 79)

# Set background
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = SOFT_IVORY

print("Processing slide 4...")

# Update all text elements
for shape in slide.shapes:
    if hasattr(shape, "text_frame"):
        tf = shape.text_frame
        text = tf.text.strip()

        # Section label "The Solution"
        if "The Solutio" in text or text == "The Solution":
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "The Solution"
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(9)
            run.font.bold = True
            run.font.color.rgb = PHTHALO_GREEN
            shape.left = Inches(0.96)
            shape.top = Inches(0.28)

        # Slide number
        elif text == "4":
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "4"
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(9)
            run.font.bold = False
            run.font.color.rgb = MUSHROOM_TAUPE
            p.alignment = PP_ALIGN.RIGHT

        # Title
        elif "Fawn is a Therapy Alternative" in text:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Fawn is a Therapy Alternative for Non-Clinical Everyday Support"
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(32)
            run.font.bold = False
            run.font.color.rgb = PHTHALO_GREEN
            shape.left = Inches(0.96)
            shape.top = Inches(0.70)
            shape.width = Inches(11.0)
            shape.height = Inches(1.0)

        # Subtitle
        elif "Miniature Robot Companions" in text or "Pixar-level Backstories" in text:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Miniature Robot Companions with Pixar-level Backstories"
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(16)
            run.font.bold = False
            run.font.color.rgb = MUSHROOM_TAUPE
            shape.left = Inches(0.96)
            shape.top = Inches(1.85)

        # Feature 1 label
        elif text == "Relational AI Agent" or "Relational AI Agent" in text:
            # Check if this is the label (not part of a longer text)
            if len(text) < 30:
                tf.clear()
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = "Relational AI Agent"
                run.font.name = "Haas Grot Text Trial"
                run.font.size = Pt(15)
                run.font.bold = False
                run.font.color.rgb = PHTHALO_GREEN

        # Feature 2 label
        elif text == "Emotionally Mature Robot" or "Emotionally Mature Robot" in text:
            if len(text) < 30:
                tf.clear()
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = "Emotionally Mature Robot"
                run.font.name = "Haas Grot Text Trial"
                run.font.size = Pt(15)
                run.font.bold = False
                run.font.color.rgb = PHTHALO_GREEN

        # Feature 3 label
        elif text == "AI-Native World Lore" or "AI-Native World Lore" in text:
            if len(text) < 30:
                tf.clear()
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = "AI-Native World Lore"
                run.font.name = "Haas Grot Text Trial"
                run.font.size = Pt(15)
                run.font.bold = False
                run.font.color.rgb = PHTHALO_GREEN

        # Quote 1: Bluebell
        elif "Bluebell is sweeter" in text:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "\u201cBluebell is sweeter than other AI I\u2019ve come across. I feel more comforted talking to her.\u201d"
            run.font.name = "Canela Deck Regular Trial"
            run.font.size = Pt(13)
            run.font.bold = False
            run.font.italic = True
            run.font.color.rgb = BODY_GRAY

        # Attribution 1: Noah
        elif "Noah" in text and "25-year-old" in text:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "- Noah, 25-year-old man"
            run.font.name = "Canela Deck Regular Trial"
            run.font.size = Pt(11)
            run.font.bold = False
            run.font.italic = True
            run.font.color.rgb = BODY_GRAY

        # Description 1
        elif "A supportive friend who" in text or "always available" in text:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "A supportive friend who\u2019s always available."
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(11)
            run.font.bold = False
            run.font.color.rgb = BODY_GRAY

        # Quote 2: Victoria
        elif "Having her nearby makes me feel" in text:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "\u201cHaving her nearby makes me feel less alone.\u201d"
            run.font.name = "Canela Deck Regular Trial"
            run.font.size = Pt(13)
            run.font.bold = False
            run.font.italic = True
            run.font.color.rgb = BODY_GRAY

        # Attribution 2: Victoria
        elif "Victoria" in text and "28-year-old" in text:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "- Victoria, 28-year-old woman"
            run.font.name = "Canela Deck Regular Trial"
            run.font.size = Pt(11)
            run.font.bold = False
            run.font.italic = True
            run.font.color.rgb = BODY_GRAY

        # Description 2
        elif "Embodiment takes pen-pal" in text or "relationships to the next level" in text:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Embodiment takes pen-pal relationships to the next level."
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(11)
            run.font.bold = False
            run.font.color.rgb = BODY_GRAY

        # Quote 3: Clare
        elif "The story reminds me so much" in text:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "\u201cThe story reminds me so much of the human creation story\u2026 I think that\u2019s why I have such a close connection to Hazel.\u201d"
            run.font.name = "Canela Deck Regular Trial"
            run.font.size = Pt(13)
            run.font.bold = False
            run.font.italic = True
            run.font.color.rgb = BODY_GRAY

        # Attribution 3: Clare
        elif "Clare" in text and "38-year-old" in text:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "- Clare, 38-year-old woman"
            run.font.name = "Canela Deck Regular Trial"
            run.font.size = Pt(11)
            run.font.bold = False
            run.font.italic = True
            run.font.color.rgb = BODY_GRAY

        # Description 3
        elif "Films, audio, news, and stories" in text:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Films, audio, news, and stories that bring meaning to the relationship."
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(11)
            run.font.bold = False
            run.font.color.rgb = BODY_GRAY

        # Numbered circles - keep as is, just ensure font
        elif text in ["1", "2", "3"] and shape.width < Inches(1):
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(18)
            run.font.bold = True
            run.font.color.rgb = PHTHALO_GREEN

print("Slide 4 updated!")
prs.save('workers/worker_A.pptx')
print("Saved to workers/worker_A.pptx")
