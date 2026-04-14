"""Edit slide 3 - Therapy vs Emotional Needs comparison"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Load presentation
prs = Presentation('workers/worker_A.pptx')
slide = prs.slides[2]  # Slide 3 (0-indexed)

# Brand colors
SOFT_IVORY = RGBColor(246, 241, 233)
PHTHALO_GREEN = RGBColor(18, 60, 51)
PHTHALO_GREEN_ALT = RGBColor(15, 77, 63)  # Alternate deeper green
WARM_BEIGE = RGBColor(231, 216, 199)
MUSHROOM_TAUPE = RGBColor(184, 169, 153)
BODY_GRAY = RGBColor(86, 83, 79)

# Set background
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = SOFT_IVORY

print("Processing slide 3...")

# Update all text elements
for shape in slide.shapes:
    if hasattr(shape, "text_frame"):
        tf = shape.text_frame
        text = tf.text.strip()

        # Section label "The Problem"
        if text == "The Problem":
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "The Problem"
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(9)
            run.font.bold = True
            run.font.color.rgb = PHTHALO_GREEN
            shape.left = Inches(0.96)
            shape.top = Inches(0.28)

        # Slide number
        elif text == "3":
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "3"
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(9)
            run.font.bold = False
            run.font.color.rgb = MUSHROOM_TAUPE
            p.alignment = PP_ALIGN.RIGHT

        # Title
        elif "Therapy Isn't Accessible" in text:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Therapy Isn't Accessible Enough To Address The Gap"
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(36)
            run.font.bold = False
            run.font.color.rgb = PHTHALO_GREEN
            shape.left = Inches(0.96)
            shape.top = Inches(1.0)

        # Left panel header "THERAPY IS"
        elif text == "THERAPY IS":
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "THERAPY IS"
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.color.rgb = MUSHROOM_TAUPE

        # Left panel items - gray text
        elif text in ["Hard to access", "Inconsistent", "Expensive", "Infrequent"]:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(21)
            run.font.bold = False
            run.font.color.rgb = BODY_GRAY

        # Right panel header "EMOTIONAL NEEDS ARE"
        elif text == "EMOTIONAL NEEDS ARE":
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "EMOTIONAL NEEDS ARE"
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.color.rgb = MUSHROOM_TAUPE

        # Right panel items - bold dark green
        elif text in ["Constant", "Unpredictable", "Immediate", "Everyday"]:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(26)
            run.font.bold = True
            run.font.color.rgb = PHTHALO_GREEN

        # Quote text - use Canela
        elif "I don't have anyone to talk to" in text:
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = "\u201cI don\u2019t have anyone to talk to when I\u2019m experiencing anxiety in the middle of the night.\u201d"
            run.font.name = "Canela Deck Bold Trial"
            run.font.size = Pt(18)
            run.font.bold = False  # Bold is baked into the font name
            run.font.color.rgb = PHTHALO_GREEN

        # Attribution - use Canela italic
        elif "Sika" in text or "Occupational Therapist" in text:
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = "— Sika, Occupational Therapist"
            run.font.name = "Canela Deck Regular Trial"
            run.font.size = Pt(14)
            run.font.bold = False
            run.font.italic = True
            run.font.color.rgb = PHTHALO_GREEN_ALT

print("Slide 3 updated!")
prs.save('workers/worker_A.pptx')
print("Saved to workers/worker_A.pptx")
