"""Edit slide 2 - Mental health gap statistics"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Load presentation
prs = Presentation('workers/worker_A.pptx')
slide = prs.slides[1]  # Slide 2 (0-indexed)

# Brand colors
SOFT_IVORY = RGBColor(246, 241, 233)
PHTHALO_GREEN = RGBColor(18, 60, 51)
WARM_BEIGE = RGBColor(231, 216, 199)
MUSHROOM_TAUPE = RGBColor(184, 169, 153)
BODY_GRAY = RGBColor(86, 83, 79)

# Set background
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = SOFT_IVORY

print("Processing slide 2...")

# Update all text elements
for shape in slide.shapes:
    if hasattr(shape, "text_frame"):
        tf = shape.text_frame
        text = tf.text.strip()

        # Section label "The Problem"
        if text == "The Problem":
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "The Problem"
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(9)
            run.font.bold = True
            run.font.color.rgb = PHTHALO_GREEN
            shape.left = Inches(0.96)
            shape.top = Inches(0.28)

        # Slide number
        elif text == "2":
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "2"
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(9)
            run.font.bold = False
            run.font.color.rgb = MUSHROOM_TAUPE
            p.alignment = PP_ALIGN.RIGHT

        # Title
        elif "Half of Americans Who Need Mental Health Support" in text:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Half of Americans Who Need Mental Health Support Don't Receive It"
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(40)
            run.font.bold = False
            run.font.color.rgb = PHTHALO_GREEN
            shape.left = Inches(0.96)
            shape.top = Inches(1.0)
            shape.width = Inches(11.5)

        # Bar 1: 60M+ Americans Need Support
        elif "60M+" in text or "Americans Need Support" in text:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "60M+ Americans Need Support"
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(18)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            # Full width bar
            shape.left = Inches(0.96)
            shape.top = Inches(2.8)
            shape.width = Inches(11.0)
            shape.height = Inches(0.6)
            # Background
            shape.fill.solid()
            shape.fill.fore_color.rgb = WARM_BEIGE

        # Bar 2: 30M Receive Care
        elif "30M" in text and "Receive Care" in text:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "30M Receive Care"
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(18)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            # Half width bar (50% of 60M)
            shape.left = Inches(0.96)
            shape.top = Inches(3.6)
            shape.width = Inches(5.5)
            shape.height = Inches(0.6)
            # Background
            shape.fill.solid()
            shape.fill.fore_color.rgb = WARM_BEIGE

        # Hero number "30M"
        elif text == "30M" and shape.width < Inches(2):
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "30M"
            run.font.name = "Haas Grot Disp Trial"
            run.font.size = Pt(72)
            run.font.bold = True
            run.font.color.rgb = PHTHALO_GREEN
            p.alignment = PP_ALIGN.CENTER
            shape.left = Inches(8.0)
            shape.top = Inches(2.8)

        # Hero label "AMERICANS WITHOUT CARE"
        elif "AMERICANS" in text and "WITHOUT" in text:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "AMERICANS\nWITHOUT\nCARE"
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = MUSHROOM_TAUPE
            p.alignment = PP_ALIGN.CENTER
            shape.left = Inches(8.0)
            shape.top = Inches(3.6)

        # Context stat
        elif "1 in 5" in text:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "1 in 5 "
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = PHTHALO_GREEN
            run2 = p.add_run()
            run2.text = "U.S. adults experience mental illness each year"
            run2.font.name = "Haas Grot Text Trial"
            run2.font.size = Pt(14)
            run2.font.bold = False
            run2.font.color.rgb = BODY_GRAY
            shape.left = Inches(0.96)
            shape.top = Inches(4.6)

        # Source attribution
        elif "National Institute" in text or "SAMHSA" in text:
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "National Institute of Mental Health; SAMHSA"
            run.font.name = "Haas Grot Text Trial"
            run.font.size = Pt(10)
            run.font.bold = False
            run.font.color.rgb = MUSHROOM_TAUPE
            shape.top = Inches(6.8)

print("Slide 2 updated!")
prs.save('workers/worker_A.pptx')
print("Saved to workers/worker_A.pptx")
