#!/usr/bin/env python3
"""
Fix broken image references in deck_final.pptx
"""

from pptx import Presentation
from pptx.util import Inches
import os
from pathlib import Path

def find_all_images():
    """Find all available images in the workspace"""
    images = {}

    # Search in images directory
    for root, dirs, files in os.walk('/workspace/images'):
        for file in files:
            if file.lower().endswith(('.png', '.jpg', '.jpeg')):
                full_path = os.path.join(root, file)
                images[file.lower()] = full_path

    return images

def check_and_fix_images(pptx_path, output_path=None):
    """Check and fix broken image references"""

    if output_path is None:
        output_path = pptx_path

    print(f"Loading presentation: {pptx_path}")
    prs = Presentation(pptx_path)

    available_images = find_all_images()
    print(f"\nFound {len(available_images)} available images")

    broken_images = []
    fixed_images = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        print(f"\n--- Slide {slide_num} ---")

        for shape_idx, shape in enumerate(slide.shapes):
            # Check if shape is a picture
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    # Try to access the image
                    image = shape.image
                    image_bytes = image.blob

                    # Get the image filename from the relationship
                    if hasattr(shape, '_element') and hasattr(shape._element, 'blipFill'):
                        blip = shape._element.blipFill.blip
                        if blip is not None:
                            embed_id = blip.embed
                            image_part = slide.part.related_part(embed_id)
                            image_filename = os.path.basename(image_part.partname)

                            print(f"  Shape {shape_idx}: {image_filename} - OK")
                    else:
                        print(f"  Shape {shape_idx}: Picture (no blip info)")

                except Exception as e:
                    print(f"  Shape {shape_idx}: BROKEN - {str(e)}")
                    broken_images.append({
                        'slide': slide_num,
                        'shape_idx': shape_idx,
                        'shape': shape,
                        'error': str(e)
                    })

    print(f"\n{'='*60}")
    print(f"Summary:")
    print(f"  Broken images: {len(broken_images)}")
    print(f"  Fixed images: {len(fixed_images)}")

    if broken_images:
        print("\nBroken image details:")
        for item in broken_images:
            print(f"  Slide {item['slide']}, Shape {item['shape_idx']}: {item['error']}")

    return {
        'broken': broken_images,
        'fixed': fixed_images,
        'available': available_images
    }

if __name__ == '__main__':
    result = check_and_fix_images('/workspace/output/deck_final.pptx')

    if result['broken']:
        print("\n" + "="*60)
        print("BROKEN IMAGES DETECTED")
        print("="*60)
        print(f"\nTotal broken: {len(result['broken'])}")
        print(f"\nAvailable images that could be used:")
        for name, path in sorted(result['available'].items())[:20]:
            print(f"  {name}: {path}")
    else:
        print("\n✓ No broken images found!")
