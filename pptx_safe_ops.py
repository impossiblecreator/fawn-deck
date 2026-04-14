"""
Safe PPTX Operations — backup, validate, and auto-rollback pipeline.

Every merge or slide edit is wrapped in a safety pipeline:
  1. Backup the target file before any modification
  2. Perform the operation
  3. Validate the result for corruption
  4. Auto-rollback to backup if validation fails

Usage:
    from pptx_safe_ops import safe_merge, safe_edit_slide, validate_deck

    # Merge with automatic corruption detection
    result = safe_merge('output/deck_final.pptx', worker_path, slide_nums=[5, 8])

    # Edit a slide with safety wrapper
    def my_edit(slide, prs):
        tf = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tf.text_frame.text = "Hello"

    result = safe_edit_slide('workers/worker_A.pptx', slide_num=5, edit_fn=my_edit)

    # Standalone validation
    report = validate_deck('output/deck_final.pptx')
    if not report.ok:
        for err in report.errors:
            print(err)
"""

import copy
import datetime
import hashlib
import os
import shutil
import zipfile
from dataclasses import dataclass, field
from lxml import etree
from pptx import Presentation
from pptx.util import Emu

WORKSPACE = os.path.dirname(os.path.abspath(__file__))

# Relationship type constants
IMAGE_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
HYPERLINK_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink"
SLIDE_LAYOUT_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
NOTES_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"
SLIDE_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"

BACKUP_DIR = os.path.join(WORKSPACE, "backups", "safe_ops")


# ---------------------------------------------------------------------------
# Validation report
# ---------------------------------------------------------------------------

@dataclass
class ValidationReport:
    """Structured result from validate_deck()."""
    path: str
    ok: bool = True
    slide_count: int = 0
    errors: list = field(default_factory=list)
    warnings: list = field(default_factory=list)

    def summary(self) -> str:
        if self.ok:
            return f"OK — {self.slide_count} slides, no issues."
        return (
            f"FAILED — {self.slide_count} slides, "
            f"{len(self.errors)} error(s), {len(self.warnings)} warning(s)."
        )


# ---------------------------------------------------------------------------
# validate_deck()
# ---------------------------------------------------------------------------

def validate_deck(pptx_path: str, expected_slide_count: int | None = None) -> ValidationReport:
    """Validate a .pptx file for corruption and structural issues.

    Checks:
      1. File opens with python-pptx without error
      2. All XML files in the ZIP parse correctly
      3. All relationship targets exist in the ZIP (no broken refs)
      4. All images referenced in slides exist in ppt/media/
      5. Slide count matches expected count (if provided)
      6. No orphaned hyperlink rels (internal hyperlinks pointing to missing targets)
      7. All text boxes have non-zero dimensions
      8. Every slide has exactly one slide-layout relationship
      9. Every media file has a registered content type
     10. Notes back-references are consistent

    Returns a ValidationReport with .ok, .errors, .warnings, and .slide_count.
    """
    report = ValidationReport(path=pptx_path)

    if not os.path.exists(pptx_path):
        report.ok = False
        report.errors.append(f"File not found: {pptx_path}")
        return report

    # --- Check 1: python-pptx can open it ---
    try:
        prs = Presentation(pptx_path)
        report.slide_count = len(prs.slides)
    except Exception as e:
        report.ok = False
        report.errors.append(f"python-pptx cannot open file: {e}")
        return report

    # --- Check 5: Slide count ---
    if expected_slide_count is not None and report.slide_count != expected_slide_count:
        report.errors.append(
            f"Slide count mismatch: expected {expected_slide_count}, got {report.slide_count}"
        )

    # --- Check 7: Text boxes with non-zero dimensions ---
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                w = shape.width if shape.width is not None else 0
                h = shape.height if shape.height is not None else 0
                if w == 0 or h == 0:
                    report.errors.append(
                        f"Slide {i}: text box '{shape.name}' has zero dimensions "
                        f"(w={Emu(w).inches:.2f}\", h={Emu(h).inches:.2f}\")"
                    )

    # --- ZIP-level checks ---
    try:
        with zipfile.ZipFile(pptx_path, "r") as z:
            all_names = set(z.namelist())

            # Check 2: Parse every XML file
            for name in sorted(all_names):
                if not name.endswith(".xml") and not name.endswith(".rels"):
                    continue
                try:
                    etree.fromstring(z.read(name))
                except etree.XMLSyntaxError as e:
                    report.errors.append(f"XML parse error in {name}: {e}")

            # Build content type registry
            registered_exts = set()
            registered_overrides = set()
            if "[Content_Types].xml" in all_names:
                ct_root = etree.fromstring(z.read("[Content_Types].xml"))
                for el in ct_root:
                    if el.tag == f"{{{CT_NS}}}Default":
                        registered_exts.add(el.get("Extension", "").lower())
                    elif el.tag == f"{{{CT_NS}}}Override":
                        pn = el.get("PartName", "").lstrip("/")
                        registered_overrides.add(pn)
            else:
                report.errors.append("[Content_Types].xml is missing")

            # Check 3, 4, 6, 8: Relationship checks
            all_rels = [n for n in all_names if n.endswith(".rels")]
            for rels_path in sorted(all_rels):
                if "_rels/" in rels_path:
                    part_dir = rels_path[:rels_path.index("_rels/")].rstrip("/")
                else:
                    part_dir = ""

                try:
                    rels_root = etree.fromstring(z.read(rels_path))
                except etree.XMLSyntaxError:
                    continue  # already caught above

                layout_count = 0
                is_slide_rels = "ppt/slides/_rels/" in rels_path

                for rel in rels_root:
                    rel_type = rel.get("Type", "")
                    target = rel.get("Target", "")
                    target_mode = rel.get("TargetMode", "Internal")
                    rid = rel.get("Id", "?")

                    if SLIDE_LAYOUT_TYPE in rel_type:
                        layout_count += 1

                    # Check 6: Orphaned hyperlink rels (internal only)
                    if HYPERLINK_REL_TYPE in rel_type and target_mode == "Internal":
                        base = (part_dir + "/") if part_dir else ""
                        resolved = os.path.normpath(base + target).replace("\\", "/")
                        if resolved not in all_names:
                            report.errors.append(
                                f"{rels_path}: orphaned internal hyperlink "
                                f"rId={rid} → {resolved} (missing)"
                            )

                    # Check 3: All relationship targets exist
                    if target_mode == "External":
                        continue

                    base = (part_dir + "/") if part_dir else ""
                    resolved = os.path.normpath(base + target).replace("\\", "/")

                    if resolved not in all_names:
                        report.errors.append(
                            f"{rels_path}: broken ref rId={rid} "
                            f"type=…/{rel_type.rsplit('/', 1)[-1]} → {resolved} (missing)"
                        )

                    # Check 4: Image refs point to existing media
                    if IMAGE_REL_TYPE in rel_type:
                        if resolved not in all_names:
                            report.errors.append(
                                f"{rels_path}: image rId={rid} → {resolved} "
                                f"not found in media folder"
                            )

                # Check 8: Slide layout count
                if is_slide_rels:
                    slide_label = rels_path.split("/")[-1].replace(".rels", "")
                    if layout_count == 0:
                        report.errors.append(f"{slide_label}: no slide-layout relationship")
                    elif layout_count > 1:
                        report.warnings.append(
                            f"{slide_label}: {layout_count} layout rels (expected 1)"
                        )

            # Check 9: Media content types
            for media_path in sorted(n for n in all_names if n.startswith("ppt/media/")):
                ext = media_path.rsplit(".", 1)[-1].lower() if "." in media_path else ""
                if ext not in registered_exts and media_path not in registered_overrides:
                    report.errors.append(f"No content type registered for {media_path}")

            # Check 10: Notes back-references
            for rels_path in sorted(
                n for n in all_names
                if n.startswith("ppt/slides/_rels/") and n.endswith(".rels")
            ):
                try:
                    rels_root = etree.fromstring(z.read(rels_path))
                except etree.XMLSyntaxError:
                    continue
                slide_name = rels_path.split("/")[-1].replace(".rels", "")
                slide_full = f"ppt/slides/{slide_name}"
                for rel in rels_root:
                    if NOTES_REL_TYPE not in rel.get("Type", ""):
                        continue
                    notes_target = rel.get("Target", "")
                    notes_zip = os.path.normpath(
                        f"ppt/slides/{notes_target}"
                    ).replace("\\", "/")
                    notes_rels_zip = notes_zip.replace(
                        "ppt/notesSlides/", "ppt/notesSlides/_rels/"
                    ) + ".rels"
                    if notes_rels_zip not in all_names:
                        continue
                    try:
                        notes_rels = etree.fromstring(z.read(notes_rels_zip))
                    except etree.XMLSyntaxError:
                        continue
                    for nrel in notes_rels:
                        ntype = nrel.get("Type", "")
                        if ("slide" in ntype
                                and "notesSlide" not in ntype
                                and "slideLayout" not in ntype):
                            back = os.path.normpath(
                                f"ppt/notesSlides/{nrel.get('Target', '')}"
                            ).replace("\\", "/")
                            if back != slide_full:
                                report.errors.append(
                                    f"{slide_name}: notes ref → {notes_zip} "
                                    f"but notes points back to {back}"
                                )

    except zipfile.BadZipFile as e:
        report.errors.append(f"Invalid ZIP file: {e}")

    report.ok = len(report.errors) == 0
    return report


# ---------------------------------------------------------------------------
# Backup / rollback
# ---------------------------------------------------------------------------

def _backup_path(pptx_path: str) -> str:
    """Generate a timestamped backup path for a .pptx file."""
    os.makedirs(BACKUP_DIR, exist_ok=True)
    basename = os.path.basename(pptx_path).replace(".pptx", "")
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S_%f")
    return os.path.join(BACKUP_DIR, f"{basename}_{timestamp}.pptx")


def create_backup(pptx_path: str) -> str:
    """Create a backup of a .pptx file. Returns the backup path."""
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"Cannot backup: {pptx_path} does not exist")
    backup = _backup_path(pptx_path)
    shutil.copy2(pptx_path, backup)
    return backup


def auto_rollback(pptx_path: str, backup_path: str, report: ValidationReport) -> bool:
    """Restore a .pptx from backup if validation failed.

    Args:
        pptx_path:   The file that was modified.
        backup_path: The pre-modification backup.
        report:      The ValidationReport from validate_deck().

    Returns:
        True if rollback was performed, False if file was clean (no rollback needed).
    """
    if report.ok:
        # Validation passed — clean up backup
        if os.path.exists(backup_path):
            os.remove(backup_path)
        return False

    # Validation failed — restore from backup
    if not os.path.exists(backup_path):
        raise FileNotFoundError(
            f"Rollback failed: backup not found at {backup_path}. "
            f"The corrupted file remains at {pptx_path}."
        )

    shutil.copy2(backup_path, pptx_path)
    # Keep the backup around for forensics — don't delete it
    return True


# ---------------------------------------------------------------------------
# safe_merge()
# ---------------------------------------------------------------------------

@dataclass
class SafeOpResult:
    """Result of a safe_merge() or safe_edit_slide() operation."""
    success: bool
    report: ValidationReport
    rolled_back: bool = False
    backup_path: str | None = None
    message: str = ""


def safe_merge(
    output_path: str,
    worker_path: str,
    slide_nums: list[int],
) -> SafeOpResult:
    """Merge slides from a worker file into the output deck with full safety pipeline.

    Pipeline:
      1. Backup output_path
      2. Copy slide XML from worker → output for each slide_num
      3. Fix media references (images)
      4. Validate result
      5. Auto-rollback if validation fails

    Args:
        output_path: Path to the output .pptx (will be modified in-place).
        worker_path: Path to the worker .pptx to pull slides from.
        slide_nums:  List of 1-indexed slide numbers to merge.

    Returns:
        SafeOpResult with success status, validation report, and rollback info.
    """
    # Import the media fix function from slide_manager
    from slide_manager import _fix_slide_media

    if not os.path.exists(output_path):
        return SafeOpResult(
            success=False,
            report=ValidationReport(path=output_path, ok=False,
                                    errors=[f"Output file not found: {output_path}"]),
            message=f"Output file not found: {output_path}",
        )

    if not os.path.exists(worker_path):
        return SafeOpResult(
            success=False,
            report=ValidationReport(path=output_path, ok=False,
                                    errors=[f"Worker file not found: {worker_path}"]),
            message=f"Worker file not found: {worker_path}",
        )

    # Step 1: Backup
    backup = create_backup(output_path)

    try:
        # Step 2: Merge slide XML
        output_prs = Presentation(output_path)
        worker_prs = Presentation(worker_path)
        expected_count = len(output_prs.slides)
        merged_slides = []

        for slide_num in slide_nums:
            idx = slide_num - 1
            if idx >= len(worker_prs.slides) or idx >= len(output_prs.slides):
                return _rollback_with_error(
                    output_path, backup, expected_count,
                    f"Slide {slide_num} out of range "
                    f"(output has {len(output_prs.slides)}, worker has {len(worker_prs.slides)})"
                )

            src_slide = worker_prs.slides[idx]
            dst_slide = output_prs.slides[idx]

            dst_slide._element.clear()
            for child in src_slide._element:
                dst_slide._element.append(copy.deepcopy(child))
            for attr, val in src_slide._element.attrib.items():
                dst_slide._element.set(attr, val)

            merged_slides.append((idx, worker_path))

        output_prs.save(output_path)

        # Step 3: Fix media references
        if merged_slides:
            _fix_slide_media(output_path, merged_slides)

        # Step 4: Validate
        report = validate_deck(output_path, expected_slide_count=expected_count)

        # Step 5: Auto-rollback if needed
        rolled_back = auto_rollback(output_path, backup, report)

        if rolled_back:
            return SafeOpResult(
                success=False,
                report=report,
                rolled_back=True,
                backup_path=backup,
                message=(
                    f"Merge of slides {slide_nums} FAILED validation — "
                    f"auto-rolled back. Errors: {report.errors}"
                ),
            )

        return SafeOpResult(
            success=True,
            report=report,
            rolled_back=False,
            message=f"Merged slides {slide_nums} from {os.path.basename(worker_path)} — validated OK.",
        )

    except Exception as e:
        # Any unexpected error — rollback
        if os.path.exists(backup):
            shutil.copy2(backup, output_path)
        return SafeOpResult(
            success=False,
            report=ValidationReport(path=output_path, ok=False, errors=[str(e)]),
            rolled_back=True,
            backup_path=backup,
            message=f"Merge failed with exception: {e}",
        )


# ---------------------------------------------------------------------------
# safe_edit_slide()
# ---------------------------------------------------------------------------

def safe_edit_slide(
    pptx_path: str,
    slide_num: int,
    edit_fn: callable,
) -> SafeOpResult:
    """Edit a single slide with full safety pipeline.

    Pipeline:
      1. Backup the .pptx
      2. Open with python-pptx, call edit_fn(slide, prs)
      3. Save
      4. Validate
      5. Auto-rollback if validation fails

    Args:
        pptx_path: Path to the .pptx file to edit.
        slide_num: 1-indexed slide number to edit.
        edit_fn:   Callable(slide, prs) that modifies the slide in-place.
                   Receives the python-pptx Slide and Presentation objects.

    Returns:
        SafeOpResult with success status, validation report, and rollback info.
    """
    if not os.path.exists(pptx_path):
        return SafeOpResult(
            success=False,
            report=ValidationReport(path=pptx_path, ok=False,
                                    errors=[f"File not found: {pptx_path}"]),
            message=f"File not found: {pptx_path}",
        )

    # Step 1: Backup
    backup = create_backup(pptx_path)

    try:
        # Step 2: Open and edit
        prs = Presentation(pptx_path)
        expected_count = len(prs.slides)
        idx = slide_num - 1

        if idx < 0 or idx >= len(prs.slides):
            # Clean up backup, return error
            os.remove(backup)
            return SafeOpResult(
                success=False,
                report=ValidationReport(
                    path=pptx_path, ok=False,
                    errors=[f"Slide {slide_num} out of range (deck has {len(prs.slides)} slides)"],
                ),
                message=f"Slide {slide_num} out of range.",
            )

        slide = prs.slides[idx]
        edit_fn(slide, prs)

        # Step 3: Save
        prs.save(pptx_path)

        # Step 4: Validate
        report = validate_deck(pptx_path, expected_slide_count=expected_count)

        # Step 5: Auto-rollback if needed
        rolled_back = auto_rollback(pptx_path, backup, report)

        if rolled_back:
            return SafeOpResult(
                success=False,
                report=report,
                rolled_back=True,
                backup_path=backup,
                message=(
                    f"Edit of slide {slide_num} FAILED validation — "
                    f"auto-rolled back. Errors: {report.errors}"
                ),
            )

        return SafeOpResult(
            success=True,
            report=report,
            rolled_back=False,
            message=f"Edited slide {slide_num} — validated OK.",
        )

    except Exception as e:
        if os.path.exists(backup):
            shutil.copy2(backup, pptx_path)
        return SafeOpResult(
            success=False,
            report=ValidationReport(path=pptx_path, ok=False, errors=[str(e)]),
            rolled_back=True,
            backup_path=backup,
            message=f"Edit failed with exception: {e}",
        )


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _rollback_with_error(
    output_path: str, backup: str, expected_count: int, error_msg: str
) -> SafeOpResult:
    """Restore from backup and return a failure result."""
    shutil.copy2(backup, output_path)
    return SafeOpResult(
        success=False,
        report=ValidationReport(
            path=output_path, ok=False,
            slide_count=expected_count, errors=[error_msg],
        ),
        rolled_back=True,
        backup_path=backup,
        message=error_msg,
    )
