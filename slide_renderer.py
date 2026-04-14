"""
Slide Renderer — renders PPTX slides to PNG.

Uses LibreOffice + pdftoppm for pixel-perfect rendering.
Falls back to Pillow if LibreOffice is unavailable.

Usage:
    from slide_renderer import render_slide, render_slides
    render_slide('deck.pptx', slide_num=5, output='slide5.png')
    render_slides('deck.pptx', slides=[4, 5], output_dir='/workspace/renders/')
"""

import subprocess
import os
import shutil
import tempfile
from pathlib import Path


def safe_font_color(font):
    """Return a human-readable string for a python-pptx font color.

    Handles all three color states without raising AttributeError:
      - Inherited (no explicit color set) → "inherited"
      - Explicit RGB                       → "#RRGGBB"
      - Theme color slot                   → "theme:<slot_name>"

    Usage in diagnostic scripts:
        from slide_renderer import safe_font_color
        for r in para.runs:
            print(safe_font_color(r.font))
    """
    try:
        ct = font.color.type
    except Exception:
        return "inherited"
    if ct is None:
        return "inherited"
    try:
        return f"#{font.color.rgb}"
    except AttributeError:
        pass
    try:
        return f"theme:{font.color.theme_color}"
    except Exception:
        return f"unknown_type:{ct}"


def _has_tool(name):
    return shutil.which(name) is not None


def _pdf_page_for_slide(pptx_path, slide_num):
    """Return the PDF page number that corresponds to a 1-indexed PPTX slide number.

    LibreOffice skips hidden slides (show="0") when converting to PDF, so the
    PDF page number can differ from the PPTX slide index.
    """
    from pptx import Presentation
    prs = Presentation(pptx_path)
    pdf_page = 0
    for i, slide in enumerate(prs.slides):
        is_hidden = slide._element.get('show', '1') == '0'
        if not is_hidden:
            pdf_page += 1
        if i + 1 == slide_num:
            if is_hidden:
                raise RuntimeError(f"Slide {slide_num} is hidden and will not appear in the PDF")
            return pdf_page
    raise RuntimeError(f"Slide {slide_num} is out of range for {pptx_path}")


def _render_with_libreoffice(pptx_path, slide_num, output_path):
    """Convert PPTX → PDF → PNG for a specific slide."""
    # Account for hidden slides: LibreOffice omits them from the PDF
    pdf_page = _pdf_page_for_slide(pptx_path, slide_num)

    with tempfile.TemporaryDirectory() as tmpdir:
        # Step 1: PPTX → PDF via LibreOffice
        # --norestore prevents LibreOffice from blocking on a crash-recovery dialog.
        # LibreOffice only allows one instance at a time, so if another worker is
        # rendering concurrently we wait and retry rather than killing their process.
        import time
        max_attempts = 6
        wait_seconds = 10
        result = None
        for attempt in range(1, max_attempts + 1):
            # Check if LibreOffice is running (skip if pgrep unavailable)
            if _has_tool("pgrep"):
                check = subprocess.run(["pgrep", "-x", "soffice"], capture_output=True)
                if check.returncode == 0:
                    print(f"  LibreOffice busy (another worker rendering) — waiting {wait_seconds}s... (attempt {attempt}/{max_attempts})")
                    time.sleep(wait_seconds)
                    continue
            result = subprocess.run(
                ["soffice", "--headless", "--norestore", "--convert-to", "pdf",
                 "--outdir", tmpdir, os.path.abspath(pptx_path)],
                capture_output=True, text=True, timeout=120
            )
            break
        else:
            raise RuntimeError(
                f"LibreOffice is still busy after {max_attempts} attempts (~{max_attempts * wait_seconds}s).\n"
                "Another worker may be stuck. Check with: pgrep soffice"
            )

        pdf_name = Path(pptx_path).stem + ".pdf"
        pdf_path = os.path.join(tmpdir, pdf_name)

        if not os.path.exists(pdf_path):
            raise RuntimeError(f"LibreOffice produced no PDF. stderr: {result.stderr}")

        # Step 2: Extract the correct page as PNG via pdftoppm
        out_prefix = os.path.join(tmpdir, "slide")
        result2 = subprocess.run(
            ["pdftoppm", "-png", "-r", "150",
             "-f", str(pdf_page), "-l", str(pdf_page),
             pdf_path, out_prefix],
            capture_output=True, text=True, timeout=60
        )

        if result2.returncode != 0:
            raise RuntimeError(f"pdftoppm failed: {result2.stderr}")

        # pdftoppm outputs slide-{pagenum}.png with zero-padded page numbers
        for f in sorted(os.listdir(tmpdir)):
            if f.startswith("slide") and f.endswith(".png"):
                src = os.path.join(tmpdir, f)
                shutil.copy2(src, output_path)
                return output_path

        raise RuntimeError("pdftoppm produced no output")


def _render_with_pillow(pptx_path, slide_num, output_path):
    """Fallback: approximate render using Pillow."""
    from PIL import Image, ImageDraw, ImageFont
    from pptx import Presentation
    from pptx.util import Inches
    from io import BytesIO

    SCALE = 150
    FR = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
    FB = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
    FI = "/usr/share/fonts/truetype/dejavu/DejaVuSerif.ttf"
    DEFAULT_COLOR = (18, 60, 51)
    BG = (241, 237, 229)

    def hex2rgb(h):
        h = h.lstrip('#')
        return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

    def get_color(font):
        try:
            if font.color and font.color.type is not None:
                try: return hex2rgb(str(font.color.rgb))
                except: pass
        except: pass
        return DEFAULT_COLOR

    def get_font(f, pt):
        px = max(8, int(pt * SCALE / 72))
        path = FB if f.bold else (FI if f.italic else FR)
        try: return ImageFont.truetype(path, px)
        except: return ImageFont.load_default()

    prs = Presentation(pptx_path)
    slide = prs.slides[slide_num - 1]
    sw, sh = prs.slide_width / Inches(1), prs.slide_height / Inches(1)
    cw, ch = int(sw * SCALE), int(sh * SCALE)
    img = Image.new('RGB', (cw, ch), BG)
    draw = ImageDraw.Draw(img)

    for shape in slide.shapes:
        lp = int(shape.left / Inches(1) * SCALE)
        tp = int(shape.top / Inches(1) * SCALE)
        wp = int(shape.width / Inches(1) * SCALE)
        hp = int(shape.height / Inches(1) * SCALE)

        if shape.shape_type == 13:
            try:
                pic = Image.open(BytesIO(shape.image.blob)).convert('RGB')
                pic = pic.resize((wp, hp), Image.LANCZOS)
                img.paste(pic, (lp, tp))
            except:
                draw.rectangle([lp, tp, lp+wp, tp+hp], outline=(180,180,180), fill=(220,215,208))
        elif hasattr(shape, "text_frame"):
            yc = tp
            for para in shape.text_frame.paragraphs:
                if not para.text.strip(): yc += int(12 * SCALE / 144); continue
                if para.space_before:
                    try:
                        sb = para.space_before
                        yc += int((sb.pt if hasattr(sb,'pt') else sb/12700) * SCALE / 72)
                    except: pass
                xc, lh = lp, 0
                for run in para.runs:
                    sp = run.font.size.pt if run.font.size else 14
                    col, fnt = get_color(run.font), get_font(run.font, sp)
                    for w in run.text.split(' '):
                        t = (xc > lp and " " or "") + w
                        bb = fnt.getbbox(t)
                        tw, th = bb[2]-bb[0], bb[3]-bb[1]
                        draw.text((xc, yc), t, fill=col, font=fnt)
                        xc += tw; lh = max(lh, th+6)
                yc += lh

    img.save(output_path, quality=95)
    return output_path


def _default_output(pptx_path, slide_num):
    """Derive a tidy output path from the pptx filename.

    workers/worker_A.pptx  →  renders/worker_A/slide_5.png
    anything else          →  renders/slide_5.png
    """
    base = Path(pptx_path).stem  # e.g. "worker_A"
    workspace = Path(__file__).parent
    if base.startswith("worker_"):
        return str(workspace / "renders" / base / f"slide_{slide_num}.png")
    return str(workspace / "renders" / f"slide_{slide_num}.png")


def render_slide(pptx_path, slide_num, output=None, **kwargs):
    """Render a slide to PNG using LibreOffice.

    Output path is inferred from the pptx filename if not provided:
      workers/worker_A.pptx  →  renders/worker_A/slide_5.png
    Always overwrites — no version suffixes.

    Requires: LibreOffice (soffice) and pdftoppm.
    Install: brew install --cask libreoffice
    """
    if not _has_tool("soffice") or not _has_tool("pdftoppm"):
        missing = []
        if not _has_tool("soffice"):
            missing.append("LibreOffice (soffice) — install: brew install --cask libreoffice")
        if not _has_tool("pdftoppm"):
            missing.append("pdftoppm — install: brew install poppler")
        raise RuntimeError(
            "Rendering requires LibreOffice + pdftoppm. Missing:\n  " + "\n  ".join(missing)
        )

    if output is None:
        output = _default_output(pptx_path, slide_num)
    os.makedirs(os.path.dirname(output), exist_ok=True)

    return _render_with_libreoffice(pptx_path, slide_num, output)


def render_slides(pptx_path, slides, output_dir=None):
    """Render multiple slides into the same directory as render_slide() would use."""
    if output_dir is None:
        # Derive dir from the pptx path (same logic as _default_output)
        base = Path(pptx_path).stem
        workspace = Path(__file__).parent
        output_dir = str(workspace / "renders" / base) if base.startswith("worker_") else str(workspace / "renders")
    os.makedirs(output_dir, exist_ok=True)
    paths = []
    for n in slides:
        path = os.path.join(output_dir, f"slide_{n}.png")
        render_slide(pptx_path, n, output=path)
        paths.append(path)
    return paths


if __name__ == '__main__':
    import sys
    if len(sys.argv) < 3:
        print("Usage: python slide_renderer.py <pptx_path> <slide_num> [output.png]")
        sys.exit(1)
    path = render_slide(sys.argv[1], int(sys.argv[2]),
                        sys.argv[3] if len(sys.argv) > 3 else None)
    print(f"Rendered: {path}")
