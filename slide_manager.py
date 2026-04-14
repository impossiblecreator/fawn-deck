"""
Slide Manager — coordinates multiple workers editing a shared PPTX deck.

Commands:
    python slide_manager.py cut-slide 6 12                         # Mark slides 6 and 12 for removal (blanks content, no renumbering)
    python slide_manager.py uncut 6                                 # Undo a cut before clean is run
    python slide_manager.py clean                                   # Physically remove cut slides from output deck
    python slide_manager.py add-slide                               # Append 1 blank slide to source + all worker files
    python slide_manager.py add-slide 3                             # Append 3 blank slides
    python slide_manager.py add-slide --after 5                     # Insert 1 blank slide after slide 5 (becomes slide 6)
    python slide_manager.py add-slide 3 --after 5                   # Insert 3 blank slides after slide 5 (become slides 6-8)
    python slide_manager.py add-slide --bg 255,255,255             # Append with custom background
    python slide_manager.py new                                     # Create blank 20-slide source deck (16:9, cream bg)
    python slide_manager.py new --slides 15                         # 15 blank slides
    python slide_manager.py new --slides 20 --bg 255,255,255        # white background
    python slide_manager.py new --slides 20 --width 13.33 --height 7.5  # custom dimensions
    python slide_manager.py whoami             # Show this worker's identity and assignments
    python slide_manager.py setup              # Create worker copies
    python slide_manager.py assign A 4 5 10    # Assign slides 4,5,10 to worker A
    python slide_manager.py assign A --all     # Assign all slides to worker A
    python slide_manager.py status             # Show assignments and progress
    python slide_manager.py merge              # Merge all assigned slides from all workers
    python slide_manager.py merge 5 8 14       # Merge only slides 5, 8, 14 (uses assignments)
    python slide_manager.py merge A 14         # Merge slide 14 from Worker A (bypasses assignments)
    python slide_manager.py reset              # Reset output deck back to source (start fresh)
    python slide_manager.py promote            # Copy deck_final.pptx → deck_original.pptx
    python slide_manager.py render A           # Render all of worker A's assigned slides
    python slide_manager.py check-changes      # Show which workers have unmerged work
    python slide_manager.py validate           # Check output deck for broken references
    python slide_manager.py validate workers/worker_A.pptx  # Validate any pptx
"""

import copy
import datetime
import hashlib
import json
import os
import subprocess
import sys
import shutil
import zipfile
from pptx import Presentation
from pptx.util import Inches
from lxml import etree

WORKSPACE = os.path.dirname(os.path.abspath(__file__))
SOURCE = f"{WORKSPACE}/source/deck_original.pptx"
WORKERS_DIR = f"{WORKSPACE}/workers"
RENDERS_DIR = f"{WORKSPACE}/renders"
OUTPUT_DIR = f"{WORKSPACE}/output"
ASSIGNMENTS_FILE = f"{WORKSPACE}/assignments.json"
WORKER_NAMES = ["A", "B", "C"]
CUTS_FILE = f"{WORKSPACE}/cuts.json"
NOTICE_FILE = f"{WORKSPACE}/.coordinator_notice"


def _file_md5(path):
    h = hashlib.md5()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def _workers_with_changes():
    """Return list of worker names whose files differ from the source deck."""
    if not os.path.exists(SOURCE):
        return []
    source_hash = _file_md5(SOURCE)
    changed = []
    for name in WORKER_NAMES:
        path = f"{WORKERS_DIR}/worker_{name}.pptx"
        if os.path.exists(path) and _file_md5(path) != source_hash:
            changed.append(name)
    return changed


def _write_notice(message):
    """Write a notice file and send a macOS system notification."""
    with open(NOTICE_FILE, "w") as f:
        f.write(message)
    try:
        subprocess.run(
            ["osascript", "-e",
             f'display notification "{message}" with title "Fawn Deck — Coordinator"'],
            capture_output=True,
        )
    except Exception:
        pass


def _check_and_clear_notice():
    """If a coordinator notice is pending, print it and delete it."""
    if not os.path.exists(NOTICE_FILE):
        return
    with open(NOTICE_FILE) as f:
        msg = f.read().strip()
    os.remove(NOTICE_FILE)
    border = "!" * 50
    print()
    print(f"  {border}")
    print(f"  COORDINATOR NOTICE:")
    for line in msg.splitlines():
        print(f"    {line}")
    print(f"  {border}")
    print()


def load_cuts():
    if os.path.exists(CUTS_FILE):
        with open(CUTS_FILE) as f:
            return set(json.load(f))
    return set()


def save_cuts(cuts):
    with open(CUTS_FILE, "w") as f:
        json.dump(sorted(cuts), f)


def load_assignments():
    if os.path.exists(ASSIGNMENTS_FILE):
        with open(ASSIGNMENTS_FILE) as f:
            return json.load(f)
    return {name: {"slides": [], "status": "idle"} for name in WORKER_NAMES}


def save_assignments(assignments):
    with open(ASSIGNMENTS_FILE, "w") as f:
        json.dump(assignments, f, indent=2)


def cmd_cut_slide(slide_nums):
    """Mark slides as cut: blank their content in all files and record them in cuts.json.

    Slides stay in place (no renumbering) so other workers' assignments are unaffected.
    Run 'clean' after the round to physically remove them from the output deck.
    Run 'uncut <nums>' to restore a slide before it's cleaned.
    """
    from pptx.dml.color import RGBColor

    cuts = load_cuts()
    cuts.update(slide_nums)
    save_cuts(cuts)

    files = [SOURCE] + [
        f"{WORKERS_DIR}/worker_{name}.pptx"
        for name in WORKER_NAMES
        if os.path.exists(f"{WORKERS_DIR}/worker_{name}.pptx")
    ]

    for filepath in files:
        if not os.path.exists(filepath):
            continue
        prs = Presentation(filepath)
        for slide_num in slide_nums:
            idx = slide_num - 1
            if idx < 0 or idx >= len(prs.slides):
                print(f"  WARNING: Slide {slide_num} out of range in {os.path.basename(filepath)}, skipping")
                continue
            slide = prs.slides[idx]
            # Remove only actual shape elements — _spTree also holds required
            # container elements (nvGrpSpPr, grpSpPr) that must NOT be removed.
            sp_tree = slide.shapes._spTree
            for shape in list(slide.shapes):
                sp_tree.remove(shape._element)
            # Set background to a pale red as a visual "cut" marker
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 220, 220)
        prs.save(filepath)
        label = os.path.relpath(filepath, WORKSPACE)
        print(f"  Blanked {label}")

    print(f"\nMarked as cut: {sorted(slide_nums)}")
    print(f"These slides are still in place. Run 'python3 slide_manager.py clean' after merging to remove them.")


def cmd_uncut_slide(slide_nums):
    """Remove slide numbers from the cut list (undo a cut before clean is run).

    This does NOT restore the slide content — it only removes the cut marker.
    The slide will remain blank; restore content manually if needed.
    """
    cuts = load_cuts()
    restored = sorted(set(slide_nums) & cuts)
    cuts -= set(slide_nums)
    save_cuts(cuts)
    if restored:
        print(f"  Removed from cut list: {restored}")
        print(f"  Note: slide content is still blank — restore manually if needed.")
    else:
        print(f"  None of {slide_nums} were in the cut list.")


def cmd_clean():
    """Physically remove all cut slides from output/deck_final.pptx.

    Only affects the output deck — source and worker files are unchanged.
    Run this after merging, before promote.
    Clears cuts.json once done.
    """
    cuts = load_cuts()
    if not cuts:
        print("  No slides marked for cutting.")
        return

    output_path = f"{OUTPUT_DIR}/deck_final.pptx"
    if not os.path.exists(output_path):
        print(f"  ERROR: {output_path} not found. Run merge first.")
        return

    # python-pptx has no delete_slide API — manipulate XML directly
    prs = Presentation(output_path)
    total = len(prs.slides)
    to_remove = sorted([n for n in cuts if 1 <= n <= total], reverse=True)

    xml_slides = prs.slides._sldIdLst
    for slide_num in to_remove:
        idx = slide_num - 1
        slide = prs.slides[idx]
        rId = prs.slides._sldIdLst[idx].get("r:id") or prs.slides._sldIdLst[idx].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        # Remove from slide list
        xml_slides.remove(xml_slides[idx])
        # Remove relationship
        if rId:
            prs.part.drop_rel(rId)

    prs.save(output_path)
    save_cuts(set())

    remaining = len(prs.slides)
    print(f"  Removed {len(to_remove)} slide(s) from output/deck_final.pptx ({total} → {remaining} slides)")
    print(f"  cuts.json cleared.")
    print(f"\nReview the deck, then run 'python3 slide_manager.py promote'.")


def cmd_add_slide(count=1, bg_rgb=None, after=None):
    """Append or insert blank slide(s) to the source deck and all worker files.

    When --after N is given, slides are inserted after position N in every file
    and assignments/cuts are renumbered so existing references stay correct.
    Without --after, slides are appended at the end (no renumbering needed).

    count:   number of slides to add (default 1)
    bg_rgb:  background colour tuple, defaults to matching the last slide's background
    after:   insert after this slide number (0 = beginning, None = append at end)
    """
    from pptx.dml.color import RGBColor

    # Block if any workers have unsaved changes (files differ from source)
    changed = _workers_with_changes()
    if changed:
        names = ", ".join(f"Worker {n}" for n in changed)
        print(f"  ERROR: Work is in progress ({names} have unsaved changes).")
        print(f"  Merge their slides and run 'promote' + 'setup' first, then try again.")
        return

    files = [SOURCE] + [
        f"{WORKERS_DIR}/worker_{name}.pptx"
        for name in WORKER_NAMES
        if os.path.exists(f"{WORKERS_DIR}/worker_{name}.pptx")
    ]
    if os.path.exists(f"{OUTPUT_DIR}/deck_final.pptx"):
        files.append(f"{OUTPUT_DIR}/deck_final.pptx")

    # Validate --after against the source deck
    if after is not None:
        source_prs = Presentation(SOURCE)
        total_before = len(source_prs.slides)
        if after < 0 or after > total_before:
            print(f"  ERROR: --after {after} is out of range. Deck has {total_before} slides (use 0–{total_before}).")
            return

    new_slide_nums = []

    for filepath in files:
        if not os.path.exists(filepath):
            continue
        prs = Presentation(filepath)
        # Use blank layout (index 6) if available, otherwise last layout
        layout_idx = min(6, len(prs.slide_layouts) - 1)
        blank_layout = prs.slide_layouts[layout_idx]

        # Detect background from last slide if not specified
        fill_rgb = bg_rgb
        if fill_rgb is None:
            try:
                last_fill = prs.slides[-1].background.fill
                if last_fill.type is not None:
                    fill_rgb = (
                        last_fill.fore_color.rgb.red,
                        last_fill.fore_color.rgb.green,
                        last_fill.fore_color.rgb.blue,
                    )
            except Exception:
                pass
        if fill_rgb is None:
            fill_rgb = (241, 237, 229)  # Fawn cream default

        for _ in range(count):
            slide = prs.slides.add_slide(blank_layout)
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*fill_rgb)

        # Move appended slides to the target position
        if after is not None and after < len(prs.slides) - count:
            xml_slides = prs.slides._sldIdLst
            new_elements = []
            for _ in range(count):
                el = xml_slides[-1]
                xml_slides.remove(el)
                new_elements.append(el)
            new_elements.reverse()
            for i, el in enumerate(new_elements):
                xml_slides.insert(after + i, el)

        total = len(prs.slides)
        if filepath == SOURCE:
            if after is not None:
                new_slide_nums = list(range(after + 1, after + count + 1))
            else:
                new_slide_nums = list(range(total - count + 1, total + 1))

        prs.save(filepath)
        label = os.path.relpath(filepath, WORKSPACE)
        print(f"  Updated {label} → now {total} slides")

    # Renumber assignments and cuts when inserting (not appending)
    if after is not None:
        assignments = load_assignments()
        changed = False
        for name in assignments:
            old_slides = assignments[name]["slides"]
            new_slides = sorted(s + count if s > after else s for s in old_slides)
            if new_slides != old_slides:
                assignments[name]["slides"] = new_slides
                changed = True
        if changed:
            save_assignments(assignments)
            print(f"  Renumbered assignments (slides > {after} shifted by +{count})")

        cuts = load_cuts()
        if cuts:
            new_cuts = {s + count if s > after else s for s in cuts}
            if new_cuts != cuts:
                save_cuts(new_cuts)
                print(f"  Renumbered cuts (slides > {after} shifted by +{count})")

    print(f"\nNew slide(s): {new_slide_nums}")
    print(f"Assign with: python3 slide_manager.py assign <worker> {' '.join(str(n) for n in new_slide_nums)}")


def cmd_new(num_slides=20, bg_rgb=(241, 237, 229), width_in=13.33, height_in=7.5):
    """Create a blank source deck and save it as source/deck_original.pptx.

    Defaults:
      --slides 20          number of blank slides
      --bg 241,237,229     background colour (R,G,B) — Fawn cream
      --width  13.33       slide width in inches  (16:9 widescreen)
      --height  7.5        slide height in inches (16:9 widescreen)

    If a source deck already exists it is backed up before being replaced.
    Run 'setup' afterwards to create worker files.
    """
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width  = Inches(width_in)
    prs.slide_height = Inches(height_in)

    blank_layout = prs.slide_layouts[6]  # index 6 = fully blank

    for _ in range(num_slides):
        slide = prs.slides.add_slide(blank_layout)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*bg_rgb)

    os.makedirs(os.path.dirname(SOURCE), exist_ok=True)

    if os.path.exists(SOURCE):
        backup = f"{WORKSPACE}/source/deck_original.bak.pptx"
        shutil.copy2(SOURCE, backup)
        print(f"  Backed up existing source → source/deck_original.bak.pptx")

    prs.save(SOURCE)
    print(f"  Created blank deck: {num_slides} slides, {width_in}\" × {height_in}\", bg=rgb{bg_rgb}")
    print(f"\nRun 'python3 slide_manager.py setup' to create worker files.")


def cmd_setup():
    """Create worker copies from the source deck.

    If any worker files already exist, they are backed up to
    workers/backups/YYYYMMDD_HHMMSS/ before being overwritten.
    """
    if not os.path.exists(SOURCE):
        print(f"ERROR: Source deck not found at {SOURCE}")
        print("Copy your deck to source/deck_original.pptx first.")
        return

    # Warn if any worker files differ from source (possible unmerged work)
    changed = _workers_with_changes()
    if changed:
        print()
        print("  WARNING — workers with unmerged changes detected:")
        for name in changed:
            print(f"    Worker {name}: workers/worker_{name}.pptx differs from source")
        print()
        print("  Setup will OVERWRITE these files (a backup will be made).")
        print("  Merge their slides first if you want to keep their work:")
        for name in changed:
            assignments = load_assignments()
            slides = assignments.get(name, {}).get("slides", [])
            if slides:
                print(f"    python3 slide_manager.py merge {' '.join(str(s) for s in slides)}  # Worker {name}")
            else:
                print(f"    python3 slide_manager.py merge {name} --all  # Worker {name} (check assignments)")
        print()
        answer = input("  Continue with setup anyway? [y/N] ").strip().lower()
        if answer != "y":
            print("  Aborted. No files were changed.")
            return
        print()

    # Back up any existing worker files before overwriting
    existing = [
        f"{WORKERS_DIR}/worker_{name}.pptx"
        for name in WORKER_NAMES
        if os.path.exists(f"{WORKERS_DIR}/worker_{name}.pptx")
    ]
    if existing:
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        backup_dir = f"{WORKERS_DIR}/backups/{timestamp}"
        os.makedirs(backup_dir, exist_ok=True)
        for worker_path in existing:
            dest = os.path.join(backup_dir, os.path.basename(worker_path))
            shutil.copy2(worker_path, dest)
        print(f"  Backed up existing worker files → workers/backups/{timestamp}/")

    prs = Presentation(SOURCE)
    total = len(prs.slides)
    print(f"Source deck: {total} slides")

    for name in WORKER_NAMES:
        worker_path = f"{WORKERS_DIR}/worker_{name}.pptx"
        shutil.copy2(SOURCE, worker_path)
        print(f"  Created worker_{name}.pptx")

    # Initialize assignments
    assignments = {name: {"slides": [], "status": "idle"} for name in WORKER_NAMES}
    save_assignments(assignments)
    print(f"\nReady. Assign slides with: python slide_manager.py assign A 4 5 10")

    # Notify workers that their files have been refreshed from source
    notice = (
        "Coordinator ran 'setup' — your worker file has been refreshed from source.\n"
        "Any unmerged local changes have been backed up but are NOT in your active file.\n"
        "Check with the coordinator before resuming work."
    )
    _write_notice(notice)
    print(f"  Notice sent to worker terminals (visible on their next command).")


def cmd_assign(worker, slide_nums):
    """Assign slides to a worker."""
    assignments = load_assignments()

    if worker not in assignments:
        print(f"ERROR: Unknown worker '{worker}'. Use one of: {WORKER_NAMES}")
        return

    # Check for conflicts
    for num in slide_nums:
        for other_worker, info in assignments.items():
            if other_worker != worker and num in info["slides"]:
                print(f"ERROR: Slide {num} is already assigned to Worker {other_worker}")
                return

    assignments[worker]["slides"] = sorted(set(assignments[worker]["slides"] + slide_nums))
    assignments[worker]["status"] = "assigned"
    save_assignments(assignments)

    print(f"Worker {worker}: assigned slides {assignments[worker]['slides']}")


def cmd_unassign(worker, slide_nums):
    """Remove slide assignments from a worker."""
    assignments = load_assignments()
    if worker not in assignments:
        print(f"ERROR: Unknown worker '{worker}'")
        return

    for num in slide_nums:
        if num in assignments[worker]["slides"]:
            assignments[worker]["slides"].remove(num)

    save_assignments(assignments)
    print(f"Worker {worker}: now has slides {assignments[worker]['slides']}")


def cmd_status():
    """Show current assignments and progress."""
    _check_and_clear_notice()
    assignments = load_assignments()
    prs = Presentation(SOURCE)
    total = len(prs.slides)

    print(f"\nDeck: {total} slides")
    print(f"{'='*50}")

    all_assigned = set()
    for name in WORKER_NAMES:
        info = assignments.get(name, {"slides": [], "status": "idle"})
        slides = info["slides"]
        status = info["status"]
        worker_file = f"{WORKERS_DIR}/worker_{name}.pptx"
        exists = "YES" if os.path.exists(worker_file) else "NO"
        print(f"\n  Worker {name}: [{status}] file={exists}")
        print(f"    Slides: {slides if slides else '(none assigned)'}")
        all_assigned.update(slides)

    unassigned = sorted(set(range(1, total + 1)) - all_assigned)
    if unassigned:
        # Show in ranges for readability
        ranges = []
        start = unassigned[0]
        end = start
        for n in unassigned[1:]:
            if n == end + 1:
                end = n
            else:
                ranges.append(f"{start}-{end}" if start != end else str(start))
                start = end = n
        ranges.append(f"{start}-{end}" if start != end else str(start))
        print(f"\n  Unassigned: {', '.join(ranges)}")
    else:
        print(f"\n  All slides assigned!")

    cuts = load_cuts()
    if cuts:
        print(f"\n  Pending cuts (will be removed on 'clean'): {sorted(cuts)}")


IMAGE_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
)
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
CT_NS  = "http://schemas.openxmlformats.org/package/2006/content-types"

MEDIA_CONTENT_TYPES = {
    "png":  "image/png",
    "jpg":  "image/jpeg",
    "jpeg": "image/jpeg",
    "gif":  "image/gif",
    "svg":  "image/svg+xml",
    "webp": "image/webp",
    "tiff": "image/tiff",
    "bmp":  "image/bmp",
}


def _slide_zip_paths(pptx_path):
    """Return list of (slide_zip_path, rels_zip_path) for each slide in order."""
    prs = Presentation(pptx_path)
    result = []
    for slide in prs.slides:
        partname = str(slide.part.partname).lstrip("/")   # ppt/slides/slide1.xml
        base, filename = partname.rsplit("/", 1)
        rels_path = f"{base}/_rels/{filename}.rels"
        result.append((partname, rels_path))
    return result


def _fix_slide_media(output_path, merged_slides):
    """Second-pass fix: copy image blobs from worker files into output ZIP.

    merged_slides: list of (slide_idx, worker_pptx_path) for every slide that
                   was merged in the first (XML) pass.

    PPTX stores images as separate files in ppt/media/. The first pass copies
    slide XML (including rId references) but not the media files themselves.
    This function patches the output ZIP so each rId points to the correct image.

    Safety: backs up the file before rewriting and restores on failure.
    """
    # Backup the python-pptx-saved file before we touch the ZIP
    backup = output_path + ".prebak"
    shutil.copy2(output_path, backup)

    try:
        _do_fix_slide_media(output_path, merged_slides)
    except Exception as e:
        print(f"  WARNING: Image-fix pass failed ({e}) — keeping python-pptx version (images may be broken)")
        shutil.copy2(backup, output_path)
    finally:
        if os.path.exists(backup):
            os.remove(backup)


def _do_fix_slide_media(output_path, merged_slides):
    """Internal implementation of the ZIP-level image fix.

    Strategy: use the WORKER'S .rels as the basis (since slide XML uses worker
    rIds), but preserve notes-slide relationships from the OUTPUT'S .rels.
    Notes and layout relationships are not referenced by rId in slide XML,
    so they must stay consistent with the output's notes parts — replacing
    them with the worker's version causes circular notes mismatches.
    """
    import tempfile

    out_slide_paths = _slide_zip_paths(output_path)
    wrk_slide_paths_cache = {}
    for _, wrk_path in merged_slides:
        if wrk_path not in wrk_slide_paths_cache:
            wrk_slide_paths_cache[wrk_path] = _slide_zip_paths(wrk_path)

    with zipfile.ZipFile(output_path, "r") as z:
        out_namelist = z.namelist()
        out_files = {name: z.read(name) for name in out_namelist}

    changed = False
    new_media_exts = set()

    for slide_idx, wrk_path in merged_slides:
        out_slide_zip, out_rels_zip = out_slide_paths[slide_idx]
        wrk_slide_zip, wrk_rels_zip = wrk_slide_paths_cache[wrk_path][slide_idx]
        wrk_slide_base = wrk_slide_zip.rsplit("/", 1)[0]

        with zipfile.ZipFile(wrk_path, "r") as wrk_zip:
            wrk_names = set(wrk_zip.namelist())
            if wrk_rels_zip not in wrk_names:
                continue

            # Use the worker's .rels as the basis (slide XML uses worker rIds),
            # but preserve notes-slide relationships from the output's .rels.
            wrk_rels_root = etree.fromstring(wrk_zip.read(wrk_rels_zip))

            # Save the output's notes relationship before we replace the .rels
            NOTES_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"
            out_notes_rel = None
            if out_rels_zip in out_files:
                out_rels_root = etree.fromstring(out_files[out_rels_zip])
                for rel_el in out_rels_root:
                    if NOTES_REL_TYPE in rel_el.get("Type", ""):
                        out_notes_rel = (rel_el.get("Id"), rel_el.get("Target"))
                        break

            # Patch image targets: copy blob into output ZIP, update Target path
            for rel_el in wrk_rels_root:
                if IMAGE_REL_TYPE not in rel_el.get("Type", ""):
                    continue

                wrk_target = rel_el.get("Target")
                wrk_media_zip = os.path.normpath(
                    os.path.join(wrk_slide_base, wrk_target)
                ).replace("\\", "/")

                if wrk_media_zip not in wrk_names:
                    continue

                img_bytes = wrk_zip.read(wrk_media_zip)
                ext = wrk_media_zip.rsplit(".", 1)[-1].lower() if "." in wrk_media_zip else "bin"

                content_hash = hashlib.md5(img_bytes).hexdigest()[:12]
                new_filename = f"image_m_{content_hash}.{ext}"
                new_zip_path = f"ppt/media/{new_filename}"

                out_files[new_zip_path] = img_bytes
                rel_el.set("Target", f"../media/{new_filename}")
                new_media_exts.add((ext, MEDIA_CONTENT_TYPES.get(ext, f"image/{ext}")))

            # Fix notes relationship: remove the worker's notes ref and
            # restore the output's original notes ref (if it had one).
            for rel_el in list(wrk_rels_root):
                if NOTES_REL_TYPE in rel_el.get("Type", ""):
                    wrk_rels_root.remove(rel_el)

            if out_notes_rel is not None:
                rid, target = out_notes_rel
                # Avoid rId conflicts with existing entries
                existing_rids = {el.get("Id") for el in wrk_rels_root}
                if rid in existing_rids:
                    # Pick a new rId that doesn't conflict
                    n = 100
                    while f"rId{n}" in existing_rids:
                        n += 1
                    rid = f"rId{n}"
                notes_el = etree.SubElement(wrk_rels_root, "Relationship")
                notes_el.set("Id", rid)
                notes_el.set("Type", NOTES_REL_TYPE)
                notes_el.set("Target", target)

            # Write the patched worker rels as the output slide's rels
            out_files[out_rels_zip] = etree.tostring(
                wrk_rels_root,
                xml_declaration=True,
                encoding="UTF-8",
                standalone=True,
            )
            changed = True

    if not changed:
        return

    # Register any new image extensions in [Content_Types].xml
    if new_media_exts and "[Content_Types].xml" in out_files:
        ct_root = etree.fromstring(out_files["[Content_Types].xml"])
        existing_exts = {
            el.get("Extension", "").lower()
            for el in ct_root
            if el.tag == f"{{{CT_NS}}}Default"
        }
        for ext, mime in sorted(new_media_exts):
            if ext not in existing_exts:
                new_el = etree.SubElement(ct_root, f"{{{CT_NS}}}Default")
                new_el.set("Extension", ext)
                new_el.set("ContentType", mime)
        out_files["[Content_Types].xml"] = etree.tostring(
            ct_root,
            xml_declaration=True,
            encoding="UTF-8",
            standalone=True,
        )

    tmp_fd, tmp_path = tempfile.mkstemp(suffix=".pptx", dir=os.path.dirname(output_path))
    os.close(tmp_fd)
    try:
        with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as new_zip:
            for name, data in out_files.items():
                new_zip.writestr(name, data)
        os.replace(tmp_path, output_path)
    except Exception:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)
        raise


def cmd_reset_output():
    """Reset output/deck_final.pptx back to the source deck.

    Use this to start a fresh merge round. Without this, merge always
    builds on top of the existing deck_final.pptx.
    """
    output_path = f"{OUTPUT_DIR}/deck_final.pptx"
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    shutil.copy2(SOURCE, output_path)
    print(f"  Reset output/deck_final.pptx from source.")


def cmd_merge(only_slides=None, from_worker=None):
    """Merge worker changes into the output deck.

    Builds on top of the existing output/deck_final.pptx so that multiple
    workers can each merge their own slides without overwriting each other.
    If no output deck exists yet, it is created from source.

    Use 'python3 slide_manager.py reset' to wipe the output and start fresh.

    from_worker: if set, only pull slides from this worker (ignores assignments
                 for the specified slides — useful when a slide was edited but
                 not formally assigned).
    only_slides: if set, only merge these slide numbers.
    """
    assignments = load_assignments()

    output_path = f"{OUTPUT_DIR}/deck_final.pptx"
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    # Build on top of existing output, or start from source if none exists
    if os.path.exists(output_path):
        print(f"  Building on existing output/deck_final.pptx")
    else:
        shutil.copy2(SOURCE, output_path)
        print(f"  No output deck found — starting from source")

    output_prs = Presentation(output_path)
    merged_count = 0
    merged_slides = []  # [(slide_idx, worker_pptx_path)] for media fix pass

    # If a specific worker is given, pull only_slides from that worker directly,
    # bypassing the assignment check entirely.
    if from_worker is not None:
        if only_slides is None:
            # No slide numbers given — merge all of this worker's assigned slides
            only_slides = set(assignments.get(from_worker, {}).get("slides", []))
            if not only_slides:
                print(f"  Worker {from_worker} has no assigned slides. Use 'assign' first.")
                return
        workers_to_process = [(from_worker, list(only_slides))]
    else:
        workers_to_process = []
        for name in WORKER_NAMES:
            info = assignments.get(name, {"slides": []})
            slides = info["slides"]
            if only_slides is not None:
                slides = [s for s in slides if s in only_slides]
            workers_to_process.append((name, slides))

        # Warn about any requested slides that aren't assigned to anyone
        if only_slides is not None:
            all_assigned = {
                s for name in WORKER_NAMES
                for s in assignments.get(name, {}).get("slides", [])
            }
            unassigned = sorted(only_slides - all_assigned)
            for slide_num in unassigned:
                print(
                    f"  WARNING: Slide {slide_num} is not assigned to any worker — skipped.\n"
                    f"           To merge it anyway: python3 slide_manager.py merge A {slide_num}"
                )

    for name, slides in workers_to_process:
        if not slides:
            continue

        worker_path = f"{WORKERS_DIR}/worker_{name}.pptx"
        if not os.path.exists(worker_path):
            print(f"  WARNING: worker_{name}.pptx not found, skipping")
            continue

        worker_prs = Presentation(worker_path)

        for slide_num in slides:
            idx = slide_num - 1
            if idx >= len(worker_prs.slides) or idx >= len(output_prs.slides):
                print(f"  WARNING: Slide {slide_num} out of range, skipping")
                continue

            # Pass 1: copy slide XML content
            src_slide = worker_prs.slides[idx]
            dst_slide = output_prs.slides[idx]

            dst_slide._element.clear()
            for child in src_slide._element:
                dst_slide._element.append(copy.deepcopy(child))
            for attr, val in src_slide._element.attrib.items():
                dst_slide._element.set(attr, val)

            merged_slides.append((idx, worker_path))
            print(f"  Merged slide {slide_num} from Worker {name}")
            merged_count += 1

    output_prs.save(output_path)

    # Pass 2: copy image blobs and fix relationship files in the output ZIP
    if merged_slides:
        print("  Fixing image references...")
        _fix_slide_media(output_path, merged_slides)

    if only_slides is not None:
        print(f"\nMerged {merged_count} slide(s) (filtered to {sorted(only_slides)}) → {output_path}")
    else:
        print(f"\nMerged {merged_count} slide(s) → {output_path}")

    # Auto-validate after every merge so broken references are caught immediately
    cmd_validate(output_path)


def cmd_promote():
    """Promote output/deck_final.pptx to become the new source/deck_original.pptx.

    The old source is backed up as source/deck_original.bak.pptx before overwriting.
    After promoting, run 'setup' to refresh all worker files from the new source.
    """
    final_path = f"{OUTPUT_DIR}/deck_final.pptx"
    if not os.path.exists(final_path):
        print(f"ERROR: {final_path} not found. Run 'merge' first.")
        return

    if os.path.exists(SOURCE):
        backup_path = f"{WORKSPACE}/source/deck_original.bak.pptx"
        shutil.copy2(SOURCE, backup_path)
        print(f"  Backed up old source → source/deck_original.bak.pptx")

    shutil.copy2(final_path, SOURCE)
    print(f"  Promoted output/deck_final.pptx → source/deck_original.pptx")
    print(f"\nRun 'python slide_manager.py setup' to refresh worker files from the new source.")


def cmd_whoami():
    """Print this worker's identity, file, and assignments based on $WORKER_ID."""
    _check_and_clear_notice()

    worker_id = os.environ.get("WORKER_ID", "").upper()

    print()
    print("  Worker Identity")
    print("  " + "=" * 30)

    if not worker_id:
        print("  ERROR: WORKER_ID environment variable is not set.")
        print()
        print("  Start Claude Code with your worker ID, e.g.:")
        print("    WORKER_ID=A claude --dangerously-skip-permissions")
        print("    WORKER_ID=B claude --dangerously-skip-permissions")
        print("    WORKER_ID=C claude --dangerously-skip-permissions")
        sys.exit(1)

    if worker_id not in WORKER_NAMES:
        print(f"  ERROR: WORKER_ID='{worker_id}' is not valid. Must be one of: {WORKER_NAMES}")
        sys.exit(1)

    worker_path = f"{WORKERS_DIR}/worker_{worker_id}.pptx"
    exists = "YES" if os.path.exists(worker_path) else "NO — run setup first"

    assignments = load_assignments()
    slides = assignments.get(worker_id, {}).get("slides", [])

    print(f"  Worker ID:    {worker_id}")
    print(f"  Working file: workers/worker_{worker_id}.pptx  (exists: {exists})")
    print(f"  Assigned slides: {slides if slides else '(none — ask coordinator to assign)'}")

    if slides:
        slide_args = " ".join(str(s) for s in slides)
        print()
        print("  Your merge command:")
        print(f"    python3 slide_manager.py merge {slide_args}")

    print()


def cmd_render(worker):
    """Render all slides assigned to a worker.

    Slides are saved to renders/worker_{X}/ and the folder is opened in Finder.
    Each render overwrites the previous — no version suffixes.
    """
    from slide_renderer import render_slide

    assignments = load_assignments()
    if worker not in assignments:
        print(f"ERROR: Unknown worker '{worker}'")
        return

    slides = assignments[worker]["slides"]
    if not slides:
        print(f"Worker {worker} has no slides assigned")
        return

    worker_path = f"{WORKERS_DIR}/worker_{worker}.pptx"
    render_dir = f"{RENDERS_DIR}/worker_{worker}"
    os.makedirs(render_dir, exist_ok=True)

    for num in slides:
        output = f"{render_dir}/slide_{num}.png"
        render_slide(worker_path, num, output)
        print(f"  Rendered slide {num} → {output}")

    print(f"\nOpening {render_dir} in Finder...")
    subprocess.run(["open", render_dir])


def cmd_validate(pptx_path=None):
    """Check a deck for issues that trigger PowerPoint's repair dialog.

    Checks:
      1. Every XML file in the ZIP parses without errors
      2. Every .rels file target (all parts, not just slides) exists in the ZIP
      3. Every file in ppt/media/ has a content type in [Content_Types].xml
      4. Every slide has exactly one slide-layout relationship
    """
    if pptx_path is None:
        pptx_path = f"{OUTPUT_DIR}/deck_final.pptx"

    if not os.path.exists(pptx_path):
        print(f"  ERROR: {pptx_path} not found.")
        return

    label = os.path.relpath(pptx_path, WORKSPACE) if pptx_path.startswith(WORKSPACE) else pptx_path
    print(f"\n  Validating {label} ...")

    errors = []
    warnings = []

    SLIDE_LAYOUT_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"

    with zipfile.ZipFile(pptx_path, "r") as z:
        all_names = set(z.namelist())

        # --- 1. Parse every XML file ---
        for name in sorted(all_names):
            if not name.endswith(".xml") and not name.endswith(".rels"):
                continue
            try:
                etree.fromstring(z.read(name))
            except etree.XMLSyntaxError as e:
                errors.append(f"XML parse error in {name}: {e}")

        # --- 2. Build registered content types ---
        registered_exts = set()
        registered_overrides = set()
        if "[Content_Types].xml" in all_names:
            ct_root = etree.fromstring(z.read("[Content_Types].xml"))
            for el in ct_root:
                if el.tag == f"{{{CT_NS}}}Default":
                    registered_exts.add(el.get("Extension", "").lower())
                elif el.tag == f"{{{CT_NS}}}Override":
                    pn = el.get("PartName", "").lstrip("/")
                    registered_overrides.add(pn)
        else:
            errors.append("[Content_Types].xml is missing")

        # --- 3. Check every .rels file (slides, charts, notes, etc.) ---
        all_rels = [n for n in all_names if n.endswith(".rels")]
        for rels_path in sorted(all_rels):
            # Derive the base directory of the part this rels file describes.
            # e.g. _rels/.rels                          → "" (root)
            #      ppt/_rels/presentation.xml.rels      → "ppt"
            #      ppt/slides/_rels/slide1.xml.rels     → "ppt/slides"
            if "_rels/" in rels_path:
                part_dir = rels_path[:rels_path.index("_rels/")].rstrip("/")
            else:
                part_dir = ""

            try:
                rels_root = etree.fromstring(z.read(rels_path))
            except etree.XMLSyntaxError:
                continue  # already caught above

            layout_count = 0
            is_slide_rels = "ppt/slides/_rels/" in rels_path

            for rel in rels_root:
                rel_type = rel.get("Type", "")
                target = rel.get("Target", "")
                target_mode = rel.get("TargetMode", "Internal")

                if SLIDE_LAYOUT_TYPE in rel_type:
                    layout_count += 1

                if target_mode == "External":
                    continue

                base = (part_dir + "/") if part_dir else ""
                resolved = os.path.normpath(base + target).replace("\\", "/")

                if resolved not in all_names:
                    errors.append(
                        f"{rels_path}: broken ref rId={rel.get('Id')} "
                        f"type=…/{rel_type.rsplit('/', 1)[-1]} → {resolved} (missing)"
                    )

            if is_slide_rels:
                slide_label = rels_path.split("/")[-1].replace(".rels", "")
                if layout_count == 0:
                    errors.append(f"{slide_label}: no slide-layout relationship")
                elif layout_count > 1:
                    warnings.append(f"{slide_label}: {layout_count} layout rels (expected 1)")

        # --- 4. Every media file must have a content type ---
        for media_path in sorted(n for n in all_names if n.startswith("ppt/media/")):
            ext = media_path.rsplit(".", 1)[-1].lower() if "." in media_path else ""
            if ext not in registered_exts and media_path not in registered_overrides:
                errors.append(f"No content type registered for {media_path}")

        # --- 5. Notes slide back-references must match ---
        NOTES_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"
        SLIDE_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
        for rels_path in sorted(n for n in all_names if n.startswith("ppt/slides/_rels/") and n.endswith(".rels")):
            try:
                rels_root = etree.fromstring(z.read(rels_path))
            except etree.XMLSyntaxError:
                continue
            slide_name = rels_path.split("/")[-1].replace(".rels", "")
            slide_full = f"ppt/slides/{slide_name}"
            for rel in rels_root:
                if NOTES_REL_TYPE not in rel.get("Type", ""):
                    continue
                notes_target = rel.get("Target", "")
                notes_zip = os.path.normpath(f"ppt/slides/{notes_target}").replace("\\", "/")
                notes_rels_zip = notes_zip.replace("ppt/notesSlides/", "ppt/notesSlides/_rels/") + ".rels"
                if notes_rels_zip not in all_names:
                    continue
                try:
                    notes_rels = etree.fromstring(z.read(notes_rels_zip))
                except etree.XMLSyntaxError:
                    continue
                for nrel in notes_rels:
                    ntype = nrel.get("Type", "")
                    if "slide" in ntype and "notesSlide" not in ntype and "slideLayout" not in ntype:
                        back = os.path.normpath(f"ppt/notesSlides/{nrel.get('Target', '')}").replace("\\", "/")
                        if back != slide_full:
                            errors.append(
                                f"{slide_name}: notes ref → {notes_zip} but notes points back to {back}"
                            )

    if errors:
        print(f"\n  ERRORS ({len(errors)}) — PowerPoint will repair/remove these:")
        for e in errors:
            print(f"    [ERROR] {e}")
    if warnings:
        print(f"\n  Warnings ({len(warnings)}):")
        for w in warnings:
            print(f"    [WARN]  {w}")
    if not errors and not warnings:
        print(f"  OK — no issues found.")
    print()


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(0)

    cmd = sys.argv[1].lower()

    if cmd == "add-slide":
        args = sys.argv[2:]
        count = 1
        bg_rgb = None
        after = None
        i = 0
        while i < len(args):
            if args[i] == "--bg" and i + 1 < len(args):
                bg_rgb = tuple(int(x) for x in args[i + 1].split(","))
                i += 2
            elif args[i] == "--after" and i + 1 < len(args):
                after = int(args[i + 1])
                i += 2
            elif args[i].isdigit():
                count = int(args[i]); i += 1
            else:
                print(f"Unknown option: {args[i]}"); sys.exit(1)
        cmd_add_slide(count=count, bg_rgb=bg_rgb, after=after)
    elif cmd == "new":
        # Parse optional --key value flags
        args = sys.argv[2:]
        kwargs = {}
        i = 0
        while i < len(args):
            if args[i] == "--slides" and i + 1 < len(args):
                kwargs["num_slides"] = int(args[i + 1]); i += 2
            elif args[i] == "--bg" and i + 1 < len(args):
                kwargs["bg_rgb"] = tuple(int(x) for x in args[i + 1].split(","))
                i += 2
            elif args[i] == "--width" and i + 1 < len(args):
                kwargs["width_in"] = float(args[i + 1]); i += 2
            elif args[i] == "--height" and i + 1 < len(args):
                kwargs["height_in"] = float(args[i + 1]); i += 2
            else:
                print(f"Unknown option: {args[i]}")
                sys.exit(1)
        cmd_new(**kwargs)
    elif cmd == "whoami":
        cmd_whoami()
    elif cmd == "setup":
        cmd_setup()
    elif cmd == "assign":
        if len(sys.argv) < 4:
            print("Usage: python slide_manager.py assign <worker> <slide_nums...>")
            print("       python slide_manager.py assign <worker> --all")
            sys.exit(1)
        worker = sys.argv[2].upper()
        if sys.argv[3] == "--all":
            prs = Presentation(SOURCE)
            slides = list(range(1, len(prs.slides) + 1))
        else:
            slides = [int(x) for x in sys.argv[3:]]
        cmd_assign(worker, slides)
    elif cmd == "unassign":
        worker = sys.argv[2].upper()
        slides = [int(x) for x in sys.argv[3:]]
        cmd_unassign(worker, slides)
    elif cmd == "status":
        cmd_status()
    elif cmd == "merge":
        args = sys.argv[2:]
        from_worker = None
        use_all = False
        if args and args[0].upper() in WORKER_NAMES:
            from_worker = args[0].upper()
            args = args[1:]
        if args and args[0] in ("--all", "-all", "all"):
            use_all = True
            args = args[1:]
        if use_all:
            # Merge all slides assigned to from_worker (or all workers if none specified)
            assignments = load_assignments()
            if from_worker:
                only = set(assignments.get(from_worker, {}).get("slides", []))
            else:
                only = {s for info in assignments.values() for s in info.get("slides", [])}
        else:
            only = set(int(x) for x in args) if args else None
        cmd_merge(only_slides=only, from_worker=from_worker)
    elif cmd == "cut-slide":
        if len(sys.argv) < 3:
            print("Usage: python slide_manager.py cut-slide <slide_nums...>")
            sys.exit(1)
        cmd_cut_slide([int(x) for x in sys.argv[2:]])
    elif cmd == "uncut":
        if len(sys.argv) < 3:
            print("Usage: python slide_manager.py uncut <slide_nums...>")
            sys.exit(1)
        cmd_uncut_slide([int(x) for x in sys.argv[2:]])
    elif cmd == "clean":
        cmd_clean()
    elif cmd == "reset":
        cmd_reset_output()
    elif cmd == "promote":
        cmd_promote()
    elif cmd == "render":
        if len(sys.argv) < 3:
            print("Usage: python slide_manager.py render <worker>")
            sys.exit(1)
        cmd_render(sys.argv[2].upper())
    elif cmd == "validate":
        target = sys.argv[2] if len(sys.argv) > 2 else None
        cmd_validate(target)
    elif cmd == "check-changes":
        changed = _workers_with_changes()
        if not changed:
            print("  No unmerged changes detected — all worker files match source.")
        else:
            assignments = load_assignments()
            print()
            print("  Workers with unmerged changes:")
            for name in changed:
                slides = assignments.get(name, {}).get("slides", [])
                slide_str = " ".join(str(s) for s in slides) if slides else "(none assigned)"
                print(f"    Worker {name}  assigned slides: {slide_str}")
                if slides:
                    print(f"      Merge: python3 slide_manager.py merge {slide_str}")
            print()
    else:
        print(f"Unknown command: {cmd}")
        print(__doc__)
