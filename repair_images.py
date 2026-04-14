#!/usr/bin/env python3
"""
Repair broken images in deck_final.pptx by copying from deck_original.pptx
"""

from pptx import Presentation
from copy import deepcopy
import os

def has_broken_images(slide):
    """Check if a slide has broken images"""
    broken_count = 0
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            try:
                _ = shape.image.blob
            except:
                broken_count += 1
    return broken_count

def copy_slide_with_images(source_prs, target_prs, slide_index):
    """
    Copy a slide from source to target presentation, replacing the target slide
    """
    # Note: This copies the entire slide, which preserves images
    source_slide = source_prs.slides[slide_index]
    target_slide = target_prs.slides[slide_index]

    # Get the slide layout
    blank_layout = target_prs.slide_layouts[6]  # Blank layout

    # Copy slide dimensions and elements
    # We need to delete all shapes and re-add from source
    shapes_to_delete = list(target_slide.shapes)
    for shape in reversed(shapes_to_delete):
        sp = shape.element
        sp.getparent().remove(sp)

    # Copy all shapes from source
    for shape in source_slide.shapes:
        # Create element copy
        new_element = deepcopy(shape.element)
        target_slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')

def repair_presentation(source_path, target_path, output_path):
    """
    Repair broken images in target by copying from source
    """
    print(f"Loading source: {source_path}")
    source_prs = Presentation(source_path)

    print(f"Loading target: {target_path}")
    target_prs = Presentation(target_path)

    print(f"\nChecking {len(target_prs.slides)} slides...")

    slides_to_fix = []

    for idx, slide in enumerate(target_prs.slides):
        broken = has_broken_images(slide)
        if broken > 0:
            print(f"  Slide {idx + 1}: {broken} broken images")
            slides_to_fix.append(idx)

    if not slides_to_fix:
        print("\n✓ No broken images found!")
        return

    print(f"\n{len(slides_to_fix)} slides need repair")
    print("\nCopying slides from source...")

    for idx in slides_to_fix:
        try:
            print(f"  Fixing slide {idx + 1}...")
            copy_slide_with_images(source_prs, target_prs, idx)
        except Exception as e:
            print(f"  ERROR on slide {idx + 1}: {e}")

    print(f"\nSaving repaired presentation to: {output_path}")
    target_prs.save(output_path)
    print("✓ Done!")

if __name__ == '__main__':
    repair_presentation(
        source_path='/workspace/source/deck_original.pptx',
        target_path='/workspace/output/deck_final.pptx',
        output_path='/workspace/output/deck_final_repaired.pptx'
    )
