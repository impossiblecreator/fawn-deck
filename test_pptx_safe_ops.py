"""
Tests for pptx_safe_ops — backup, validation, and auto-rollback pipeline.

Run:  python3 -m pytest test_pptx_safe_ops.py -v
"""

import os
import shutil
import tempfile
import zipfile
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor

import pytest

from pptx_safe_ops import (
    validate_deck,
    create_backup,
    auto_rollback,
    safe_edit_slide,
    safe_merge,
    ValidationReport,
)


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

def _make_deck(path, num_slides=3):
    """Create a minimal valid .pptx with N blank slides."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]  # blank layout
    for _ in range(num_slides):
        slide = prs.slides.add_slide(layout)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(241, 237, 229)
    prs.save(path)
    return path


@pytest.fixture
def tmp_dir():
    d = tempfile.mkdtemp(prefix="safe_ops_test_")
    yield d
    shutil.rmtree(d, ignore_errors=True)


@pytest.fixture
def valid_deck(tmp_dir):
    return _make_deck(os.path.join(tmp_dir, "valid.pptx"), num_slides=5)


@pytest.fixture
def worker_deck(tmp_dir):
    """Create a worker deck with text on slide 2 for merge testing."""
    path = os.path.join(tmp_dir, "worker.pptx")
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    for i in range(5):
        slide = prs.slides.add_slide(layout)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(241, 237, 229)
        if i == 1:  # slide 2
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            txBox.text_frame.text = "Worker edit on slide 2"
    prs.save(path)
    return path


# ---------------------------------------------------------------------------
# Tests: validate_deck()
# ---------------------------------------------------------------------------

class TestValidateDeck:
    """Tests for the validate_deck() function."""

    def test_valid_deck_passes(self, valid_deck):
        report = validate_deck(valid_deck)
        assert report.ok is True
        assert report.slide_count == 5
        assert len(report.errors) == 0

    def test_missing_file(self, tmp_dir):
        report = validate_deck(os.path.join(tmp_dir, "nonexistent.pptx"))
        assert report.ok is False
        assert any("not found" in e for e in report.errors)

    def test_slide_count_mismatch(self, valid_deck):
        report = validate_deck(valid_deck, expected_slide_count=10)
        assert report.ok is False
        assert any("Slide count mismatch" in e for e in report.errors)

    def test_slide_count_match(self, valid_deck):
        report = validate_deck(valid_deck, expected_slide_count=5)
        assert report.ok is True

    def test_broken_image_ref(self, valid_deck):
        """Inject a broken image relationship into a slide's .rels."""
        with zipfile.ZipFile(valid_deck, "r") as z:
            files = {name: z.read(name) for name in z.namelist()}

        # Find a slide rels file and add a broken image ref
        slide_rels = None
        for name in files:
            if name.startswith("ppt/slides/_rels/") and name.endswith(".rels"):
                slide_rels = name
                break

        assert slide_rels is not None, "No slide rels found in test deck"

        root = etree.fromstring(files[slide_rels])
        fake_rel = etree.SubElement(root, "Relationship")
        fake_rel.set("Id", "rId999")
        fake_rel.set("Type", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image")
        fake_rel.set("Target", "../media/ghost_image.png")
        files[slide_rels] = etree.tostring(root, xml_declaration=True, encoding="UTF-8")

        # Rewrite the ZIP
        with zipfile.ZipFile(valid_deck, "w", zipfile.ZIP_DEFLATED) as z:
            for name, data in files.items():
                z.writestr(name, data)

        report = validate_deck(valid_deck)
        assert report.ok is False
        assert any("ghost_image.png" in e for e in report.errors)

    def test_orphaned_hyperlink(self, valid_deck):
        """Inject an internal hyperlink pointing to a missing target."""
        with zipfile.ZipFile(valid_deck, "r") as z:
            files = {name: z.read(name) for name in z.namelist()}

        slide_rels = None
        for name in files:
            if name.startswith("ppt/slides/_rels/") and name.endswith(".rels"):
                slide_rels = name
                break

        root = etree.fromstring(files[slide_rels])
        fake_rel = etree.SubElement(root, "Relationship")
        fake_rel.set("Id", "rId998")
        fake_rel.set("Type", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink")
        fake_rel.set("Target", "slide99.xml")
        # Internal (no TargetMode="External")
        files[slide_rels] = etree.tostring(root, xml_declaration=True, encoding="UTF-8")

        with zipfile.ZipFile(valid_deck, "w", zipfile.ZIP_DEFLATED) as z:
            for name, data in files.items():
                z.writestr(name, data)

        report = validate_deck(valid_deck)
        assert report.ok is False
        assert any("orphaned internal hyperlink" in e for e in report.errors)

    def test_external_hyperlink_not_flagged(self, valid_deck):
        """External hyperlinks should NOT be flagged as orphaned."""
        with zipfile.ZipFile(valid_deck, "r") as z:
            files = {name: z.read(name) for name in z.namelist()}

        slide_rels = None
        for name in files:
            if name.startswith("ppt/slides/_rels/") and name.endswith(".rels"):
                slide_rels = name
                break

        root = etree.fromstring(files[slide_rels])
        fake_rel = etree.SubElement(root, "Relationship")
        fake_rel.set("Id", "rId997")
        fake_rel.set("Type", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink")
        fake_rel.set("Target", "https://example.com")
        fake_rel.set("TargetMode", "External")
        files[slide_rels] = etree.tostring(root, xml_declaration=True, encoding="UTF-8")

        with zipfile.ZipFile(valid_deck, "w", zipfile.ZIP_DEFLATED) as z:
            for name, data in files.items():
                z.writestr(name, data)

        report = validate_deck(valid_deck)
        assert report.ok is True

    def test_zero_dimension_textbox(self, tmp_dir):
        """A text box with zero width should be flagged."""
        path = os.path.join(tmp_dir, "zero_dim.pptx")
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        # Add a textbox with zero width
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Emu(0), Inches(1))
        txBox.text_frame.text = "I'm invisible"
        prs.save(path)

        report = validate_deck(path)
        assert report.ok is False
        assert any("zero dimensions" in e for e in report.errors)

    def test_missing_content_type_for_media(self, valid_deck):
        """A media file without a registered content type should be flagged."""
        with zipfile.ZipFile(valid_deck, "r") as z:
            files = {name: z.read(name) for name in z.namelist()}

        # Add a fake media file with an unregistered extension
        files["ppt/media/fake_file.xyz"] = b"not a real image"

        with zipfile.ZipFile(valid_deck, "w", zipfile.ZIP_DEFLATED) as z:
            for name, data in files.items():
                z.writestr(name, data)

        report = validate_deck(valid_deck)
        assert report.ok is False
        assert any("No content type" in e and "fake_file.xyz" in e for e in report.errors)

    def test_malformed_xml(self, valid_deck):
        """A file with broken XML should be caught (either by python-pptx or XML parser)."""
        with zipfile.ZipFile(valid_deck, "r") as z:
            files = {name: z.read(name) for name in z.namelist()}

        # Corrupt a non-slide XML so python-pptx can still open the file
        # but the ZIP-level XML check catches it. Use a rels file.
        rels_file = None
        for name in files:
            if name.startswith("ppt/slides/_rels/") and name.endswith(".rels"):
                rels_file = name
                break

        files[rels_file] = b"<broken xml that does not close"

        with zipfile.ZipFile(valid_deck, "w", zipfile.ZIP_DEFLATED) as z:
            for name, data in files.items():
                z.writestr(name, data)

        report = validate_deck(valid_deck)
        assert report.ok is False
        assert any("XML parse error" in e or "cannot open" in e for e in report.errors)

    def test_corrupt_zip(self, tmp_dir):
        """A non-ZIP file should be caught."""
        path = os.path.join(tmp_dir, "corrupt.pptx")
        with open(path, "wb") as f:
            f.write(b"this is not a zip file at all")

        report = validate_deck(path)
        assert report.ok is False
        assert len(report.errors) > 0


# ---------------------------------------------------------------------------
# Tests: create_backup() and auto_rollback()
# ---------------------------------------------------------------------------

class TestBackupRollback:

    def test_create_backup(self, valid_deck):
        backup = create_backup(valid_deck)
        assert os.path.exists(backup)
        # Backup content matches original
        with open(valid_deck, "rb") as f1, open(backup, "rb") as f2:
            assert f1.read() == f2.read()
        os.remove(backup)

    def test_backup_missing_file_raises(self, tmp_dir):
        with pytest.raises(FileNotFoundError):
            create_backup(os.path.join(tmp_dir, "nope.pptx"))

    def test_rollback_on_failure(self, valid_deck, tmp_dir):
        """auto_rollback restores file when report.ok is False."""
        backup = create_backup(valid_deck)
        original_size = os.path.getsize(valid_deck)

        # Corrupt the file
        with open(valid_deck, "wb") as f:
            f.write(b"corrupted")

        report = ValidationReport(path=valid_deck, ok=False, errors=["test corruption"])
        rolled_back = auto_rollback(valid_deck, backup, report)

        assert rolled_back is True
        assert os.path.getsize(valid_deck) == original_size

    def test_no_rollback_on_success(self, valid_deck):
        """auto_rollback cleans up backup when report.ok is True."""
        backup = create_backup(valid_deck)
        assert os.path.exists(backup)

        report = ValidationReport(path=valid_deck, ok=True)
        rolled_back = auto_rollback(valid_deck, backup, report)

        assert rolled_back is False
        assert not os.path.exists(backup)  # backup cleaned up

    def test_rollback_missing_backup_raises(self, valid_deck, tmp_dir):
        report = ValidationReport(path=valid_deck, ok=False, errors=["bad"])
        with pytest.raises(FileNotFoundError, match="backup not found"):
            auto_rollback(valid_deck, os.path.join(tmp_dir, "nope.pptx"), report)


# ---------------------------------------------------------------------------
# Tests: safe_edit_slide()
# ---------------------------------------------------------------------------

class TestSafeEditSlide:

    def test_successful_edit(self, valid_deck):
        """A clean edit passes validation and returns success."""
        def add_text(slide, prs):
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            txBox.text_frame.text = "Safe edit test"

        result = safe_edit_slide(valid_deck, slide_num=1, edit_fn=add_text)
        assert result.success is True
        assert result.rolled_back is False
        assert result.report.ok is True

        # Verify the text is actually there
        prs = Presentation(valid_deck)
        texts = [s.text_frame.text for s in prs.slides[0].shapes if hasattr(s, "text_frame")]
        assert "Safe edit test" in texts

    def test_edit_out_of_range(self, valid_deck):
        """Editing a slide beyond range returns failure without rollback."""
        def noop(slide, prs):
            pass

        result = safe_edit_slide(valid_deck, slide_num=99, edit_fn=noop)
        assert result.success is False
        assert any("out of range" in e for e in result.report.errors)

    def test_edit_exception_triggers_rollback(self, valid_deck):
        """An exception in edit_fn triggers rollback to pre-edit state."""
        prs_before = Presentation(valid_deck)
        count_before = len(prs_before.slides)

        def bad_edit(slide, prs):
            # Add something, then crash
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            txBox.text_frame.text = "This should be rolled back"
            raise RuntimeError("Intentional test failure")

        result = safe_edit_slide(valid_deck, slide_num=1, edit_fn=bad_edit)
        assert result.success is False
        assert result.rolled_back is True

        # File should be restored to original
        prs_after = Presentation(valid_deck)
        assert len(prs_after.slides) == count_before

    def test_edit_missing_file(self, tmp_dir):
        result = safe_edit_slide(
            os.path.join(tmp_dir, "nope.pptx"), slide_num=1, edit_fn=lambda s, p: None
        )
        assert result.success is False

    def test_edit_creating_zero_dim_textbox_fails(self, valid_deck):
        """An edit that creates a zero-dimension textbox should fail validation and rollback."""
        def bad_dimensions(slide, prs):
            slide.shapes.add_textbox(Inches(1), Inches(1), Emu(0), Inches(1))

        result = safe_edit_slide(valid_deck, slide_num=1, edit_fn=bad_dimensions)
        assert result.success is False
        assert result.rolled_back is True
        assert any("zero dimensions" in e for e in result.report.errors)


# ---------------------------------------------------------------------------
# Tests: safe_merge()
# ---------------------------------------------------------------------------

class TestSafeMerge:

    def test_successful_merge(self, valid_deck, worker_deck):
        """Merging a clean worker slide succeeds."""
        result = safe_merge(valid_deck, worker_deck, slide_nums=[2])
        assert result.success is True
        assert result.rolled_back is False
        assert result.report.ok is True

        # Verify the worker's text landed
        prs = Presentation(valid_deck)
        texts = []
        for shape in prs.slides[1].shapes:
            if hasattr(shape, "text_frame"):
                texts.append(shape.text_frame.text)
        assert "Worker edit on slide 2" in texts

    def test_merge_out_of_range_slide(self, valid_deck, worker_deck):
        result = safe_merge(valid_deck, worker_deck, slide_nums=[99])
        assert result.success is False
        assert result.rolled_back is True

    def test_merge_missing_output(self, tmp_dir, worker_deck):
        result = safe_merge(
            os.path.join(tmp_dir, "nope.pptx"), worker_deck, slide_nums=[1]
        )
        assert result.success is False

    def test_merge_missing_worker(self, valid_deck, tmp_dir):
        result = safe_merge(
            valid_deck, os.path.join(tmp_dir, "nope.pptx"), slide_nums=[1]
        )
        assert result.success is False

    def test_merge_preserves_slide_count(self, valid_deck, worker_deck):
        """Merge should not change the total slide count."""
        prs_before = Presentation(valid_deck)
        count_before = len(prs_before.slides)

        result = safe_merge(valid_deck, worker_deck, slide_nums=[1, 2, 3])
        assert result.success is True
        assert result.report.slide_count == count_before


# ---------------------------------------------------------------------------
# Tests: ValidationReport
# ---------------------------------------------------------------------------

class TestValidationReport:

    def test_summary_ok(self):
        r = ValidationReport(path="test.pptx", ok=True, slide_count=10)
        assert "OK" in r.summary()

    def test_summary_failed(self):
        r = ValidationReport(path="test.pptx", ok=False, slide_count=10,
                             errors=["err1", "err2"], warnings=["w1"])
        s = r.summary()
        assert "FAILED" in s
        assert "2 error(s)" in s
        assert "1 warning(s)" in s
