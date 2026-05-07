from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PPTX = "workers/worker_A.pptx"
SLIDE_INDEX = 2

FONT_TEXT = "Haas Grot Text Trial"

GREEN = RGBColor(0x12, 0x3C, 0x33)
SOFT_GREEN = RGBColor(0x6D, 0x84, 0x7B)
BODY = RGBColor(0x56, 0x53, 0x4F)
IVORY = RGBColor(0xF6, 0xF1, 0xE8)
BEIGE = RGBColor(0xE7, 0xD8, 0xC7)
TAUPE = RGBColor(0xB8, 0xA9, 0x99)
SAGE = RGBColor(0x9C, 0xAF, 0x9A)
HONEY = RGBColor(0xE8, 0xCF, 0xA4)


def remove_owned_shapes(slide):
    for shape in list(slide.shapes):
        if shape.name.startswith("ff_s3_"):
            shape._element.getparent().remove(shape._element)


def set_text(shape, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = FONT_TEXT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_textbox(slide, name, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    return set_text(shape, text, size, color, bold, align)


def add_step(slide, idx, x, title, body):
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x),
        Inches(3.02),
        Inches(0.56),
        Inches(0.56),
    )
    circle.name = f"ff_s3_step_{idx}_number"
    circle.fill.solid()
    circle.fill.fore_color.rgb = HONEY
    circle.line.fill.background()
    set_text(circle, str(idx), 13, GREEN, bold=True, align=PP_ALIGN.CENTER)
    circle.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_textbox(
        slide,
        f"ff_s3_step_{idx}_title",
        x - 0.10,
        3.86,
        3.02,
        0.58,
        title,
        17,
        GREEN,
        bold=True,
    )
    add_textbox(
        slide,
        f"ff_s3_step_{idx}_body",
        x - 0.10,
        4.56,
        3.20,
        0.62,
        body,
        14,
        BODY,
    )


def add_arrow(slide, name, x, y, w, h, direction="right"):
    shape_type = MSO_SHAPE.RIGHT_ARROW if direction == "right" else MSO_SHAPE.LEFT_ARROW
    arrow = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    arrow.name = name
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = SAGE
    arrow.line.fill.background()
    return arrow


def add_line(slide, name, x1, y1, x2, y2, color=SAGE, width=1.7, arrow=False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.name = name
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if arrow:
        line.line.end_arrowhead = True
    return line


def main():
    prs = Presentation(PPTX)
    slide = prs.slides[SLIDE_INDEX]

    remove_owned_shapes(slide)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = IVORY

    add_textbox(
        slide,
        "ff_s3_section_label",
        0.96,
        0.28,
        1.55,
        0.24,
        "The Solution",
        9,
        SOFT_GREEN,
        bold=True,
    )

    add_textbox(
        slide,
        "ff_s3_headline",
        0.92,
        0.82,
        11.00,
        0.78,
        "Building Social Robots For Everyone",
        34,
        GREEN,
    )

    add_line(slide, "ff_s3_main_path", 1.30, 3.30, 11.95, 3.30, width=1.8, arrow=True)

    add_step(
        slide,
        1,
        1.87,
        "AI capable of relationship",
        "Built for one archetype.",
    )
    add_step(
        slide,
        2,
        5.87,
        "Robot capable of relationship",
        "Built for that same archetype.",
    )
    add_step(
        slide,
        3,
        9.87,
        "Next archetype",
        "Use technology + money earned.",
    )

    repeat_icon = slide.shapes.add_shape(
        MSO_SHAPE.CIRCULAR_ARROW,
        Inches(9.90),
        Inches(5.62),
        Inches(0.36),
        Inches(0.36),
    )
    repeat_icon.name = "ff_s3_repeat_icon"
    repeat_icon.fill.solid()
    repeat_icon.fill.fore_color.rgb = SAGE
    repeat_icon.line.fill.background()
    add_textbox(
        slide,
        "ff_s3_repeat_label",
        10.34,
        5.66,
        1.00,
        0.28,
        "repeat",
        10,
        TAUPE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    prs.save(PPTX)


if __name__ == "__main__":
    main()
